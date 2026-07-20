"""
Custom PPTX template support — "upload your company deck, AI fills it".

Two endpoints:

POST /parse-template
    Takes a signed URL to a user-uploaded .pptx. Returns:
    - a template spec: every slide's editable text SLOTS (id, role,
      sample text, size hints) so the frontend can show what the AI
      will fill and let users manage a template library
    - an extracted BRAND KIT: theme colors mapped to our palette keys
      + heading/body fonts. Lovable can pass these as
      branding.custom_colors / branding.font to /export-pptx so even
      generator-built decks match the company brand.

POST /fill-template
    Takes the same signed URL + a query + document chunks. The AI writes
    content into the template's text slots; we rewrite the text INSIDE
    the original file, preserving every bit of design (backgrounds,
    logos, fonts, colors, positioning). Uploads the filled deck to R2.

How users should author templates: any normal .pptx works. Every text
box/placeholder that contains text becomes a fillable slot. Sample text
length guides the AI (a short title stays short, a 3-bullet list gets
~3 bullets), so templates should contain realistic dummy text.
"""
import io
import re
import copy
import json
import uuid
import httpx
import ai
from lxml import etree
from fastapi import APIRouter, HTTPException
from pydantic import BaseModel
from typing import Optional
from pptx import Presentation
from dotenv import load_dotenv

from presentation import upload_to_r2

load_dotenv()

router = APIRouter()

# Cap how much of the template's sample text we send to the AI per slot
SAMPLE_TEXT_CHARS = 200
MAX_SLOTS = 60          # sanity cap for absurd templates
DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


class ParseTemplateRequest(BaseModel):
    signed_url: str
    file_name:  Optional[str] = "template.pptx"


class FillTemplateRequest(BaseModel):
    signed_url: str                     # the stored template file
    query:      str                     # what the presentation should be about
    chunks:     list[str]               # retrieved document chunks from /query
    tone:       Optional[str] = "professional"
    title:      Optional[str] = None    # optional deck title hint


def _download(signed_url: str) -> bytes:
    try:
        res = httpx.get(signed_url, timeout=60, follow_redirects=True)
        res.raise_for_status()
        return res.content
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not download template file: {e}")


# ── Template analysis ───────────────────────────────────────────────────────────

def _slot_role(shape) -> str:
    """Best-effort role for a text shape, used to guide the AI."""
    try:
        if shape.is_placeholder:
            type_name = str(shape.placeholder_format.type)
            if "TITLE" in type_name:
                return "title"
            if "SUBTITLE" in type_name:
                return "subtitle"
            if "BODY" in type_name or "OBJECT" in type_name:
                return "body"
    except Exception:
        pass
    return "text"


def _iter_text_shapes(shapes):
    """Yields text-bearing shapes, descending into groups (shape_type 6)."""
    for shape in shapes:
        if shape.shape_type == 6:
            yield from _iter_text_shapes(shape.shapes)
        elif getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            yield shape


def analyze_template(pptx_bytes: bytes) -> dict:
    """Builds the slot spec the AI fills and the frontend displays."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = []
    total_slots = 0

    for s_idx, slide in enumerate(prs.slides):
        slots = []
        for shape in _iter_text_shapes(slide.shapes):
            if total_slots >= MAX_SLOTS:
                break
            text = shape.text_frame.text.strip()
            slots.append({
                "id":          f"s{s_idx}_id{shape.shape_id}",
                "shape_id":    shape.shape_id,
                "role":        _slot_role(shape),
                "sample_text": text[:SAMPLE_TEXT_CHARS],
                "paragraphs":  len(shape.text_frame.paragraphs),
                "chars":       len(text),
            })
            total_slots += 1

        layout_name = ""
        try:
            layout_name = slide.slide_layout.name or ""
        except Exception:
            pass

        slides.append({"index": s_idx, "layout": layout_name, "slots": slots})

    return {"slide_count": len(prs.slides), "slides": slides, "total_slots": total_slots}


def extract_brand_kit(pptx_bytes: bytes) -> dict:
    """
    Pulls theme colors + fonts out of /ppt/theme/theme1.xml and maps them
    to our palette keys. Best-effort: any missing piece falls back to
    sensible defaults rather than failing the parse.
    """
    colors: dict = {}
    fonts = {"heading": "Calibri", "body": "Calibri"}
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        theme_blob = None
        for part in prs.part.package.iter_parts():
            if str(part.partname).startswith("/ppt/theme/"):
                theme_blob = part.blob
                break
        if theme_blob:
            root = etree.fromstring(theme_blob)
            ns = {"a": DRAWINGML_NS}

            def scheme_color(name: str) -> Optional[str]:
                el = root.find(f".//a:clrScheme/a:{name}", ns)
                if el is None:
                    return None
                srgb = el.find("a:srgbClr", ns)
                if srgb is not None:
                    return srgb.get("val")
                sys = el.find("a:sysClr", ns)
                if sys is not None:
                    return sys.get("lastClr")
                return None

            raw = {name: scheme_color(name)
                   for name in ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2")}
            colors = {
                "primary":    raw.get("accent1") or "1E2761",
                "secondary":  raw.get("lt2") or "E8E8F0",
                "accent":     raw.get("accent2") or "FFFFFF",
                "dark_text":  raw.get("dk1") or raw.get("dk2") or "1A1A2E",
                "light_text": raw.get("lt1") or "F5F5F5",
            }

            major = root.find(".//a:fontScheme/a:majorFont/a:latin", ns)
            minor = root.find(".//a:fontScheme/a:minorFont/a:latin", ns)
            if major is not None and major.get("typeface"):
                fonts["heading"] = major.get("typeface")
            if minor is not None and minor.get("typeface"):
                fonts["body"] = minor.get("typeface")
    except Exception as e:
        print(f"[templates] Brand kit extraction failed (non-fatal): {e}")

    if not colors:
        colors = {"primary": "1E2761", "secondary": "E8E8F0", "accent": "FFFFFF",
                  "dark_text": "1A1A2E", "light_text": "F5F5F5"}
    return {"colors": {k: f"#{v.lstrip('#')}" for k, v in colors.items()}, "fonts": fonts}


# ── Text replacement that preserves formatting ──────────────────────────────────

def _set_paragraph_text(paragraph, text: str):
    """Replaces a paragraph's text while keeping its first run's formatting."""
    runs = paragraph.runs
    if runs:
        for run in runs[1:]:
            run._r.getparent().remove(run._r)
        runs[0].text = text
    else:
        paragraph.text = text


def replace_shape_text(text_frame, new_text: str):
    """
    Writes new_text into a text frame line-by-line ('\n' = new paragraph),
    preserving per-paragraph formatting (bullet style, font, color, size):
    - existing paragraphs keep their own formatting and just get new text
    - extra lines clone the LAST paragraph's formatting
    - surplus paragraphs are removed
    """
    lines = [line for line in str(new_text).split("\n") if line.strip()] or [""]
    paragraphs = list(text_frame.paragraphs)
    txBody = text_frame._txBody

    while len(paragraphs) < len(lines):
        clone = copy.deepcopy(paragraphs[-1]._p)
        paragraphs[-1]._p.addnext(clone)
        paragraphs = list(text_frame.paragraphs)

    for paragraph in paragraphs[len(lines):]:
        paragraph._p.getparent().remove(paragraph._p)
    paragraphs = list(text_frame.paragraphs)

    for paragraph, line in zip(paragraphs, lines):
        _set_paragraph_text(paragraph, line)


# ── Routes ──────────────────────────────────────────────────────────────────────

@router.post("/parse-template")
def parse_template(request: ParseTemplateRequest):
    """
    Analyzes an uploaded .pptx template: fillable text slots per slide +
    extracted brand kit. Lovable stores this spec alongside the file in
    its own DB to build the workspace's template library.
    """
    try:
        pptx_bytes = _download(request.signed_url)
        try:
            spec = analyze_template(pptx_bytes)
        except Exception:
            raise HTTPException(
                status_code=400,
                detail="This file could not be read as a PowerPoint (.pptx) template."
            )

        if spec["total_slots"] == 0:
            raise HTTPException(
                status_code=400,
                detail="No text found in this template. Add sample text to the text "
                       "boxes you want the AI to fill (e.g. 'Headline goes here')."
            )

        brand = extract_brand_kit(pptx_bytes)
        return {
            "success":   True,
            "file_name": request.file_name,
            **spec,
            "brand":     brand,
        }
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"PARSE-TEMPLATE ERROR: {e}")
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Template parsing failed: {e}")


# NOTE: plain `def` — blocking LLM call runs in FastAPI's threadpool.
@router.post("/fill-template")
def fill_template(request: FillTemplateRequest):
    """
    Fills a custom template with AI-written content from the workspace's
    documents. The original file's design is untouched — only text inside
    existing text boxes is replaced. Returns the filled deck's R2 URL plus
    the per-slot text (so Lovable can show a content preview).
    """
    try:
        if not request.chunks:
            raise HTTPException(
                status_code=400,
                detail="chunks cannot be empty — run /query first and pass its chunks."
            )

        pptx_bytes = _download(request.signed_url)
        try:
            prs  = Presentation(io.BytesIO(pptx_bytes))
            spec = analyze_template(pptx_bytes)
        except HTTPException:
            raise
        except Exception:
            raise HTTPException(status_code=400, detail="This file could not be read as a .pptx template.")

        if spec["total_slots"] == 0:
            raise HTTPException(status_code=400, detail="Template has no text slots to fill.")

        # ── Ask the AI to fill every slot ──────────────────────────────────────
        context = "\n\n---\n\n".join(request.chunks[:12])
        slots_description = json.dumps(
            [
                {
                    "slide": s["index"],
                    "slots": [
                        {
                            "id":          slot["id"],
                            "role":        slot["role"],
                            "sample_text": slot["sample_text"],
                            "lines":       slot["paragraphs"],
                            "max_chars":   max(60, int(slot["chars"] * 1.3)),
                        }
                        for slot in s["slots"]
                    ],
                }
                for s in spec["slides"]
            ],
            indent=1,
        )

        title_hint = f'\nDeck title hint: "{request.title}"' if request.title else ""
        system_prompt = f"""You are filling an existing PowerPoint template with new content.
Tone: {request.tone}. Audience: CEOs and senior managers.

You get the template's text slots (with the sample text currently in them) and source
document content. Write NEW text for every slot:
- Keep each slot's PURPOSE: a title slot gets a short title, a bullet-list slot gets bullets.
- Match the sample's shape: same number of lines (separate lines with \\n), and stay under
  max_chars per slot — text that is too long will overflow the design.
- Use ONLY facts from the source content. Mark estimates with "(est.)".
- Slots you cannot sensibly fill (page numbers, legal footers, dates): return the sample text unchanged.

Respond ONLY with valid JSON, no markdown fences:
{{ "fills": {{ "<slot id>": "new text", ... }} }}
Every slot id from the input MUST appear in fills."""

        user_prompt = f"""Template slots:
{slots_description}

The presentation should be about: "{request.query}"{title_hint}

Source content from company documents:
{context}

Fill every slot. Respond only with the JSON object."""

        result = ai.chat_json(
            messages=[{"role": "user", "content": user_prompt}],
            system=system_prompt,
            max_tokens=4000,
            temperature=0.4,
        )
        fills = result.get("fills") or {}
        if not isinstance(fills, dict) or not fills:
            raise HTTPException(status_code=500, detail="AI returned no slot fills.")

        # ── Write fills into the actual file ───────────────────────────────────
        filled_count = 0
        preview = []
        for s_idx, slide in enumerate(prs.slides):
            slide_preview = []
            for shape in _iter_text_shapes(slide.shapes):
                slot_id  = f"s{s_idx}_id{shape.shape_id}"
                new_text = fills.get(slot_id)
                if new_text is None or not str(new_text).strip():
                    continue
                replace_shape_text(shape.text_frame, str(new_text))
                filled_count += 1
                slide_preview.append({"id": slot_id, "text": str(new_text)})
            preview.append({"slide": s_idx, "fills": slide_preview})

        if filled_count == 0:
            raise HTTPException(status_code=500, detail="AI fills did not match any template slots.")

        buf = io.BytesIO()
        prs.save(buf)
        filled_bytes = buf.getvalue()

        safe_title = re.sub(r"[^A-Za-z0-9_-]+", "_", (request.title or request.query))[:50] or "presentation"
        key = f"presentations/{uuid.uuid4()}/{safe_title}.pptx"
        url = upload_to_r2(
            filled_bytes, key,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        print(f"[fill-template] Filled {filled_count}/{spec['total_slots']} slots -> {key}")

        return {
            "success":      True,
            "pptx_url":     url,
            "title":        request.title or request.query,
            "slide_count":  spec["slide_count"],
            "slots_total":  spec["total_slots"],
            "slots_filled": filled_count,
            "preview":      preview,
            "size_kb":      round(len(filled_bytes) / 1024, 1),
        }

    except HTTPException:
        raise
    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"AI returned invalid JSON: {e}")
    except Exception as e:
        import traceback
        print(f"FILL-TEMPLATE ERROR: {e}")
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Template filling failed: {e}")

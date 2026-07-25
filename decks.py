"""
Gamma-style deck engine ("AI Create" presentations v2).

The Gamma flow, mapped to endpoints:
  1. GET  /deck-themes           → one-click restyle themes (light + dark)
  2. POST /generate-deck-outline → prompt + chunks → editable outline
                                   (title + card list) the user approves
  3. POST /generate-card         → generates ONE card's blocks; the
                                   frontend calls this per outline card in
                                   sequence, so the deck appears card by
                                   card like Gamma's streaming generation
  4. POST /edit-card             → AI edit of a single card ("shorter",
                                   "add a chart", "more punchy")
  5. POST /export-deck           → renders cards to a real PPTX on R2

A CARD is web-native and block-based (not a fixed 16:9 layout):
  { id, title, layout: "default" | "image_right",
    blocks: [
      {type: "text",     text}
      {type: "bullets",  items: [str]}
      {type: "stats",    items: [{value, label, sublabel}]}
      {type: "chart",    chart: {type, labels, datasets}, insight}
      {type: "table",    columns: [str], rows: [[str]]}
      {type: "callout",  style: "info"|"success"|"warning", text}
      {type: "quote",    text, attribution}
      {type: "image",    query, url, credit, credit_link}
    ] }

Lovable renders cards as a scrollable document (Notion-like editing);
export maps each card to one PPTX slide.
"""
import io
import re
import json
import uuid
import httpx
import ai
from fastapi import APIRouter, HTTPException
from pydantic import BaseModel
from typing import Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from dotenv import load_dotenv

from presentation import (
    PALETTES, add_text_box, headline_fit, set_slide_background,
    add_picture_fit, upload_to_r2,
)
from course import fetch_unsplash_image_url

load_dotenv()

router = APIRouter()

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

VALID_BLOCK_TYPES = ("text", "bullets", "stats", "chart", "table", "callout", "quote", "image")
VALID_CHART_TYPES = ("column", "bar", "line", "pie")

# ── Themes: palettes + light/dark mode + font pairing ───────────────────────────
THEMES = {
    "aurora_light":   {"palette": "midnight_executive", "mode": "light",
                       "name": "Aurora Light",  "description": "Clean white cards, navy accents",
                       "fonts": {"heading": "Calibri", "body": "Calibri"}},
    "midnight_dark":  {"palette": "midnight_executive", "mode": "dark",
                       "name": "Midnight",      "description": "Deep navy cards, high contrast",
                       "fonts": {"heading": "Calibri", "body": "Calibri"}},
    "ocean_light":    {"palette": "ocean_gradient", "mode": "light",
                       "name": "Ocean",         "description": "Airy white with ocean blues",
                       "fonts": {"heading": "Calibri", "body": "Calibri"}},
    "emerald_light":  {"palette": "emerald_boardroom", "mode": "light",
                       "name": "Emerald",       "description": "Fresh greens on white",
                       "fonts": {"heading": "Calibri", "body": "Calibri"}},
    "graphite_dark":  {"palette": "graphite_gold", "mode": "dark",
                       "name": "Graphite Gold", "description": "Near-black with gold accents",
                       "fonts": {"heading": "Georgia", "body": "Calibri"}},
    "burgundy_light": {"palette": "burgundy_classic", "mode": "light",
                       "name": "Burgundy",      "description": "Classic wine tones on white",
                       "fonts": {"heading": "Georgia", "body": "Calibri"}},
}


def _theme_colors(theme_id: str) -> dict:
    """Resolved colors for a theme: card bg, text, accent — used by export.
    The frontend gets the same values from /deck-themes for CSS."""
    theme   = THEMES.get(theme_id, THEMES["aurora_light"])
    palette = PALETTES[theme["palette"]]
    if theme["mode"] == "dark":
        return {"bg": palette["primary"], "text": palette["light_text"],
                "muted": palette["secondary"], "accent": palette["accent"],
                "accent_soft": palette["secondary"]}
    return {"bg": "FFFFFF", "text": palette["dark_text"],
            "muted": "6B7280", "accent": palette["primary"],
            "accent_soft": palette["secondary"]}


class OutlineRequest(BaseModel):
    query:     str
    chunks:    list[str]
    num_cards: Optional[int] = 8
    tone:      Optional[str] = "professional"


class CardPlan(BaseModel):
    id:    Optional[str] = None
    title: str
    brief: str                       # one sentence: what this card covers


class GenerateCardRequest(BaseModel):
    deck_title:    str
    card:          CardPlan
    chunks:        list[str]
    tone:          Optional[str]  = "professional"
    enable_images: Optional[bool] = True
    # Layout control (added for Presentation Studio V2). The frontend sends the
    # layouts its chosen template actually offers; the model picks ONE and
    # writes blocks that fit it. Position lets it reserve cover/closing.
    available_layouts: Optional[list[str]] = None
    template_slant:    Optional[str]       = None
    position:          Optional[int]       = None   # 0-based index of this card
    total_cards:       Optional[int]       = None


class EditCardRequest(BaseModel):
    card:        dict
    instruction: str
    chunks:      Optional[list[str]] = []
    # An instruction like "make this a timeline" legitimately changes the
    # layout, but only to one the chosen template actually offers.
    available_layouts: Optional[list[str]] = None


class ExportDeckRequest(BaseModel):
    title:    str
    subtitle: Optional[str] = ""
    cards:    list[dict]
    theme:    Optional[str] = "aurora_light"


# ── Layout catalogue ────────────────────────────────────────────────────────────
# A layout is an ARRANGEMENT, not a slide number. The model picks one per card
# from whatever its template offers, and writes blocks that suit it. Each entry
# states what the layout is for and what block shape it needs, because the model
# choosing "stats" without emitting a stats block produces an empty slide.

LAYOUT_SPECS: dict[str, str] = {
    "cover":         "Opening slide. Just a strong title and at most ONE short text block. No bullets.",
    "section":       "Divider announcing a new part. Short title, at most one short text block.",
    "title_content": "The general-purpose slide. Title plus 2-4 blocks of any kind.",
    "text_image":    "A photo beside the words. MUST include exactly one image block plus text or bullets.",
    "two_column":    "Two balanced groups side by side. MUST include two bullets blocks.",
    "stats":         "Dominated by big numbers. MUST include a stats block (2-4 figures).",
    "timeline":      "A chronological sequence. MUST include a bullets block whose items each begin with a period or date label.",
    "comparison":    "Two options weighed against each other. MUST include a table (2-3 columns) or two bullets blocks.",
    "process":       "Ordered steps. MUST include a bullets block whose items are sequential steps.",
    "chart":         "The chart is the point. MUST include a chart block; keep other blocks to one short insight.",
    "quote":         "One memorable line. MUST include a quote block and little else.",
    "closing":       "Final slide: takeaways or next steps. Bullets or a callout work best.",
}

ALL_LAYOUTS = list(LAYOUT_SPECS.keys())

# Kept so decks saved before layouts existed still render.
LEGACY_LAYOUTS = {"default": "title_content", "image_right": "text_image"}


def _coerce_layout(value, allowed: Optional[list[str]], blocks: list[dict]) -> str:
    """Validates the model's layout choice against the template, then against
    what the card actually contains. A 'chart' layout with no chart block is a
    broken slide, so we fall back rather than trust the label."""
    allowed_set = [l for l in (allowed or ALL_LAYOUTS) if l in LAYOUT_SPECS] or ["title_content"]
    layout = LEGACY_LAYOUTS.get(str(value), str(value))

    if layout not in allowed_set:
        layout = None

    types = {b["type"] for b in blocks}
    requires = {
        "chart":      "chart",
        "stats":      "stats",
        "quote":      "quote",
        "text_image": "image",
    }
    need = requires.get(layout)
    if need and need not in types:
        layout = None                      # claimed a layout it can't fill

    if layout is None:
        # Derive from content: same precedence the frontend rules use.
        for cand, btype in (("chart", "chart"), ("stats", "stats"),
                            ("quote", "quote"), ("text_image", "image")):
            if btype in types and cand in allowed_set:
                return cand
        return "title_content" if "title_content" in allowed_set else allowed_set[0]
    return layout


# ── Normalization ───────────────────────────────────────────────────────────────

def normalize_card(card: dict, enable_images: bool = True,
                   allowed_layouts: Optional[list[str]] = None) -> dict:
    """Coerces an AI-produced card into the strict block schema the
    frontend and exporter rely on. Invalid blocks are dropped, never fatal."""
    blocks = []
    for raw in (card.get("blocks") or [])[:6]:
        btype = raw.get("type")
        if btype == "text" and str(raw.get("text", "")).strip():
            blocks.append({"type": "text", "text": str(raw["text"])})
        elif btype == "bullets":
            items = [str(i) for i in (raw.get("items") or []) if str(i).strip()][:6]
            if items:
                blocks.append({"type": "bullets", "items": items})
        elif btype == "stats":
            items = [
                {"value": str(s.get("value", "")), "label": str(s.get("label", "")),
                 "sublabel": str(s.get("sublabel", ""))}
                for s in (raw.get("items") or []) if s.get("value")
            ][:4]
            if items:
                blocks.append({"type": "stats", "items": items})
        elif btype == "chart":
            chart = raw.get("chart") or {}
            if chart.get("type") not in VALID_CHART_TYPES:
                chart["type"] = "column"
            labels = [str(l) for l in (chart.get("labels") or [])]
            datasets = []
            for ds in (chart.get("datasets") or []):
                d = list(ds.get("data") or [])
                datasets.append({"label": str(ds.get("label", "Series")),
                                 "data": (d + [None] * len(labels))[:len(labels)]})
            if labels and datasets:
                blocks.append({"type": "chart",
                               "chart": {"type": chart["type"], "labels": labels, "datasets": datasets},
                               "insight": str(raw.get("insight", ""))})
        elif btype == "table":
            cols = [str(c) for c in (raw.get("columns") or [])][:6]
            rows = [
                ([str(v) if v is not None else "" for v in r] + [""] * len(cols))[:len(cols)]
                for r in (raw.get("rows") or [])[:8]
            ]
            if cols and rows:
                blocks.append({"type": "table", "columns": cols, "rows": rows})
        elif btype == "callout" and str(raw.get("text", "")).strip():
            style = raw.get("style") if raw.get("style") in ("info", "success", "warning") else "info"
            blocks.append({"type": "callout", "style": style, "text": str(raw["text"])})
        elif btype == "quote" and str(raw.get("text", "")).strip():
            blocks.append({"type": "quote", "text": str(raw["text"]),
                           "attribution": str(raw.get("attribution", ""))})
        elif btype == "image":
            query = str(raw.get("query", "")).strip()
            if raw.get("url"):
                # already resolved (round-tripping an existing card) — keep as-is
                blocks.append({"type": "image", "query": query,
                               "url": str(raw["url"]), "thumb": str(raw.get("thumb", raw["url"])),
                               "credit": str(raw.get("credit", "")),
                               "credit_link": str(raw.get("credit_link", ""))})
            elif query and enable_images:
                img = fetch_unsplash_image_url(query)
                if img:
                    blocks.append({"type": "image", "query": query, **img})

    return {
        "id":     card.get("id") or f"card_{uuid.uuid4().hex[:8]}",
        "title":  str(card.get("title", "Card")),
        "layout": _coerce_layout(card.get("layout"), allowed_layouts, blocks),
        "blocks": blocks,
    }


# ── Routes ──────────────────────────────────────────────────────────────────────

@router.get("/deck-themes")
async def deck_themes():
    return {
        "themes": [
            {
                "id": tid,
                "name": t["name"],
                "description": t["description"],
                "mode": t["mode"],
                "fonts": t["fonts"],
                "colors": {k: f"#{v}" for k, v in _theme_colors(tid).items()},
            }
            for tid, t in THEMES.items()
        ],
        "default_theme": "aurora_light",
        "block_types": list(VALID_BLOCK_TYPES),
    }


@router.post("/generate-deck-outline")
def generate_deck_outline(request: OutlineRequest):
    """Step 1 of the Gamma flow: a short editable outline, NOT full content."""
    try:
        if not request.chunks:
            raise HTTPException(status_code=400, detail="chunks cannot be empty — run /query first.")
        num_cards = max(3, min(int(request.num_cards or 8), 16))
        context = "\n\n---\n\n".join(request.chunks[:12])

        system_prompt = f"""You are planning a presentation for CEOs and senior managers.
Tone: {request.tone}. Plan {num_cards} cards that tell one coherent story:
open with the key message, develop it with evidence, close with recommendations/next steps.

Respond ONLY with valid JSON, no markdown fences:
{{
  "title": "deck title",
  "subtitle": "one-line subtitle",
  "cards": [ {{"title": "card title — a complete insight, not a topic", "brief": "one sentence: what this card will show"}} ]
}}"""
        user_prompt = f"""The presentation is about: "{request.query}"

Source content from company documents:
{context}

Plan exactly {num_cards} cards. Respond only with the JSON object."""

        data = ai.chat_json(
            messages=[{"role": "user", "content": user_prompt}],
            system=system_prompt, max_tokens=1500, temperature=0.4,
        )
        cards = [
            {"id": f"card_{uuid.uuid4().hex[:8]}",
             "title": str(c.get("title", "Card")), "brief": str(c.get("brief", ""))}
            for c in (data.get("cards") or [])[:num_cards]
        ]
        if not cards:
            raise HTTPException(status_code=500, detail="AI returned no outline cards.")
        return {
            "title":    str(data.get("title", request.query)),
            "subtitle": str(data.get("subtitle", "")),
            "cards":    cards,
        }
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"DECK-OUTLINE ERROR: {e}"); print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Outline generation failed: {e}")


@router.post("/generate-card")
def generate_card(request: GenerateCardRequest):
    """Step 2: full content for ONE outline card. The frontend calls this
    once per card, appending each result — Gamma's progressive reveal."""
    try:
        if not request.chunks:
            raise HTTPException(status_code=400, detail="chunks cannot be empty.")
        context = "\n\n---\n\n".join(request.chunks[:12])

        # Layouts this template offers, described so the model knows what each
        # arrangement needs. It chooses one and writes blocks that fill it.
        allowed = [l for l in (request.available_layouts or ALL_LAYOUTS) if l in LAYOUT_SPECS] \
                  or ["title_content"]
        layout_menu = "\n".join(f'- "{l}": {LAYOUT_SPECS[l]}' for l in allowed)

        position_hint = ""
        if request.position is not None and request.total_cards:
            if request.position == 0 and "cover" in allowed:
                position_hint = "\nThis is the FIRST card — prefer the cover layout."
            elif request.position == request.total_cards - 1 and "closing" in allowed:
                position_hint = "\nThis is the LAST card — prefer the closing layout."
            else:
                position_hint = (f"\nThis is card {request.position + 1} of {request.total_cards} "
                                 f"— do not use cover or closing.")

        slant = f"\nTemplate style: {request.template_slant}" if request.template_slant else ""

        system_prompt = f"""You are writing ONE card of the presentation "{request.deck_title}".
Tone: {request.tone}. Audience: CEOs and senior managers. Cards are concise — 2 to 4 blocks.{slant}

FIRST choose the layout that best fits what this card has to say, then write blocks that fill it.
Pick the arrangement the CONTENT calls for, not the one that comes next in a sequence.

Available layouts:
{layout_menu}{position_hint}

The layout you choose dictates the blocks you MUST emit — a "chart" layout with no chart
block, or a "stats" layout with no stats block, produces a broken slide.

Available block types (use the best 2-4 for THIS card, vary across the deck):
- {{"type": "text", "text": "short paragraph, max 60 words"}}
- {{"type": "bullets", "items": ["max 6 punchy bullets"]}}
- {{"type": "stats", "items": [{{"value": "23%", "label": "Revenue growth", "sublabel": "vs last quarter"}}]}} (2-4 stats)
- {{"type": "chart", "chart": {{"type": "column|bar|line|pie", "labels": [...], "datasets": [{{"label": "...", "data": [...]}}]}}, "insight": "one sentence"}}
- {{"type": "table", "columns": [...], "rows": [[...]]}} (max 6 cols, 8 rows)
- {{"type": "callout", "style": "info|success|warning", "text": "one key point to highlight"}}
- {{"type": "quote", "text": "...", "attribution": "..."}}
- {{"type": "image", "query": "2-4 word stock photo search phrase"}}

Rules:
- Use ONLY facts from the source content; mark estimates "(est.)".
- Charts/stats only where the source has real numbers — never invent figures to justify a layout.
- At most one image block, only if a photo genuinely helps.
- If the content does not suit any specialised layout, use "title_content". That is
  a good answer, not a failure.

Respond ONLY with valid JSON, no markdown fences:
{{"title": "card title", "layout": "one of the layouts above", "blocks": [ ... ]}}"""

        user_prompt = f"""Card to write: "{request.card.title}"
This card should show: {request.card.brief}

Source content:
{context}

Respond only with the JSON object."""

        data = ai.chat_json(
            messages=[{"role": "user", "content": user_prompt}],
            system=system_prompt, max_tokens=1800, temperature=0.4,
        )
        data["id"] = request.card.id
        if not data.get("title"):
            data["title"] = request.card.title
        card = normalize_card(data, enable_images=request.enable_images,
                              allowed_layouts=allowed)
        if not card["blocks"]:
            card["blocks"] = [{"type": "text", "text": request.card.brief}]
        return card
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"GENERATE-CARD ERROR: {e}"); print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Card generation failed: {e}")


@router.post("/edit-card")
def edit_card(request: EditCardRequest):
    """AI edit of one card, preserving the block schema."""
    try:
        context = ""
        if request.chunks:
            context = "\n\nSource content (for factual additions):\n" + "\n---\n".join(request.chunks[:6])
        edit_allowed = [l for l in (request.available_layouts or ALL_LAYOUTS)
                        if l in LAYOUT_SPECS] or ["title_content"]
        edit_menu = "\n".join(f'- "{l}": {LAYOUT_SPECS[l]}' for l in edit_allowed)

        system_prompt = f"""You edit one presentation card based on a user instruction.
Keep the same JSON structure: {{"id", "title", "layout", "blocks": [...]}} with the same
block type vocabulary already used in the card. Keep the same id. Only change what the
instruction asks.

If the instruction implies a different arrangement ("make this a timeline", "turn these
into big numbers"), change "layout" to match AND emit the blocks that layout requires.
Otherwise keep the existing layout.

Layouts available:
{edit_menu}

Respond ONLY with the updated card JSON, no markdown fences."""
        user_prompt = f"""Current card:
{json.dumps(request.card, indent=1)}

Instruction: {request.instruction}{context}

Return the updated card JSON only."""
        data = ai.chat_json(
            messages=[{"role": "user", "content": user_prompt}],
            system=system_prompt, max_tokens=1800, temperature=0.4,
        )
        data["id"] = request.card.get("id") or data.get("id")
        if not data.get("layout"):
            data["layout"] = request.card.get("layout")
        card = normalize_card(data, allowed_layouts=edit_allowed)
        if not card["blocks"]:
            raise HTTPException(status_code=500, detail="Edit produced an empty card — try rephrasing the instruction.")
        return card
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"EDIT-CARD ERROR: {e}"); print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Card edit failed: {e}")


# ── PPTX export: one card = one slide, blocks stacked top-down ──────────────────

CALLOUT_FILL = {"info": "DBEAFE", "success": "D1FAE5", "warning": "FEF3C7"}
CALLOUT_TEXT = {"info": "1E40AF", "success": "065F46", "warning": "92400E"}


def _download_image(url: str) -> Optional[bytes]:
    try:
        res = httpx.get(url, timeout=15, follow_redirects=True)
        res.raise_for_status()
        return res.content
    except Exception as e:
        print(f"[decks] image download failed: {e}")
        return None


# ── Layout-specific slide renderers ─────────────────────────────────────────────
# Each mirrors what the web editor draws for that layout, so the exported deck
# matches what the user approved on screen. They all take a fresh slide and own
# the whole canvas — no shared title bar, because several layouts (cover,
# section, quote) deliberately do not have one.

def _blank_slide(prs, colors):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, colors["bg"])
    return slide


def _first(blocks, btype):
    return next((b for b in blocks if b["type"] == btype), None)


def _accent_rule(slide, colors, left=0.9, top=3.05, width=1.5, height=0.06):
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(colors["accent"])
    bar.line.fill.background()


def _render_cover(prs, card, colors, fonts):
    """Full-bleed opening: oversized title, accent rule, one supporting line."""
    slide = _blank_slide(prs, colors)
    title = card.get("title", "")
    size, height, _ = headline_fit(title, box_width_in=10.5, max_font=54)
    add_text_box(slide, title, Inches(0.9), Inches(2.3), Inches(10.5), Inches(max(height, 1.4)),
                 font_name=fonts["heading"], font_size=size, bold=True, color=colors["text"])
    _accent_rule(slide, colors, top=2.3 + max(height, 1.4) + 0.15)
    sub = _first(card.get("blocks") or [], "text")
    if sub:
        add_text_box(slide, sub["text"][:220], Inches(0.9), Inches(2.3 + max(height, 1.4) + 0.45),
                     Inches(9.0), Inches(1.0), font_name=fonts["body"], font_size=17,
                     color=colors["muted"])


def _render_section(prs, card, colors, fonts):
    """Divider: accent band down the left, large title, optional one-liner."""
    from pptx.enum.shapes import MSO_SHAPE
    slide = _blank_slide(prs, colors)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.55), SLIDE_H)
    band.fill.solid(); band.fill.fore_color.rgb = _rgb(colors["accent"]); band.line.fill.background()

    title = card.get("title", "")
    size, height, _ = headline_fit(title, box_width_in=10.5, max_font=42)
    add_text_box(slide, title, Inches(1.3), Inches(2.9), Inches(10.5), Inches(max(height, 1.1)),
                 font_name=fonts["heading"], font_size=size, bold=True, color=colors["text"])
    sub = _first(card.get("blocks") or [], "text")
    if sub:
        add_text_box(slide, sub["text"][:180], Inches(1.3), Inches(2.9 + max(height, 1.1) + 0.2),
                     Inches(9.5), Inches(0.9), font_name=fonts["body"], font_size=15,
                     color=colors["muted"])


def _render_quote(prs, card, colors, fonts):
    """One memorable line, set large and centred."""
    slide = _blank_slide(prs, colors)
    q = _first(card.get("blocks") or [], "quote")
    text = q["text"] if q else card.get("title", "")
    add_text_box(slide, "“", Inches(0.9), Inches(1.3), Inches(2.0), Inches(1.4),
                 font_name=fonts["heading"], font_size=90, bold=True, color=colors["accent"])
    size, height, _ = headline_fit(text, box_width_in=10.6, max_font=32)
    add_text_box(slide, text, Inches(1.35), Inches(2.5), Inches(10.6), Inches(max(height, 1.8)),
                 font_name=fonts["heading"], font_size=size, italic=True, color=colors["text"])
    if q and q.get("attribution"):
        add_text_box(slide, f"— {q['attribution']}", Inches(1.35),
                     Inches(2.5 + max(height, 1.8) + 0.25), Inches(8.0), Inches(0.5),
                     font_name=fonts["body"], font_size=15, color=colors["muted"])


def _render_stats(prs, card, colors, fonts):
    """Numbers are the message: title small, figures very large."""
    slide = _blank_slide(prs, colors)
    add_text_box(slide, card.get("title", ""), Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.9),
                 font_name=fonts["heading"], font_size=26, bold=True, color=colors["text"])
    _accent_rule(slide, colors, top=1.65, width=1.2)

    stats = (_first(card.get("blocks") or [], "stats") or {}).get("items", [])
    if stats:
        col_w = 11.5 / max(len(stats), 1)
        for i, s in enumerate(stats):
            x = 0.9 + i * col_w
            add_text_box(slide, str(s.get("value", "")), Inches(x), Inches(2.6),
                         Inches(col_w - 0.3), Inches(1.5),
                         font_name=fonts["heading"], font_size=60, bold=True, color=colors["accent"])
            add_text_box(slide, str(s.get("label", "")), Inches(x), Inches(4.15),
                         Inches(col_w - 0.3), Inches(0.5),
                         font_name=fonts["body"], font_size=15, bold=True, color=colors["text"])
            if s.get("sublabel"):
                add_text_box(slide, str(s["sublabel"]), Inches(x), Inches(4.65),
                             Inches(col_w - 0.3), Inches(0.5),
                             font_name=fonts["body"], font_size=12, color=colors["muted"])

    note = _first(card.get("blocks") or [], "text")
    if note:
        add_text_box(slide, note["text"][:200], Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.8),
                     font_name=fonts["body"], font_size=13, color=colors["muted"])


def _render_timeline(prs, card, colors, fonts):
    """Horizontal rail with a marker per step."""
    from pptx.enum.shapes import MSO_SHAPE
    slide = _blank_slide(prs, colors)
    add_text_box(slide, card.get("title", ""), Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.9),
                 font_name=fonts["heading"], font_size=26, bold=True, color=colors["text"])

    items = (_first(card.get("blocks") or [], "bullets") or {}).get("items", [])[:5]
    if not items:
        return _stack_blocks(slide, card, colors, fonts, start_y=1.9)

    rail_y = 3.4
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(rail_y),
                                  Inches(11.3), Inches(0.045))
    rail.fill.solid(); rail.fill.fore_color.rgb = _rgb(colors["muted"]); rail.line.fill.background()

    step_w = 11.3 / len(items)
    for i, item in enumerate(items):
        cx = 1.0 + i * step_w + step_w / 2
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.13), Inches(rail_y - 0.11),
                                     Inches(0.26), Inches(0.26))
        dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(colors["accent"]); dot.line.fill.background()
        # Alternate above/below so longer labels do not collide.
        top = rail_y - 1.45 if i % 2 == 0 else rail_y + 0.45
        add_text_box(slide, str(item), Inches(cx - step_w / 2 + 0.15), Inches(top),
                     Inches(step_w - 0.3), Inches(1.2),
                     font_name=fonts["body"], font_size=13, color=colors["text"])


def _render_process(prs, card, colors, fonts):
    """Numbered steps flowing left to right."""
    from pptx.enum.shapes import MSO_SHAPE
    slide = _blank_slide(prs, colors)
    add_text_box(slide, card.get("title", ""), Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.9),
                 font_name=fonts["heading"], font_size=26, bold=True, color=colors["text"])

    items = (_first(card.get("blocks") or [], "bullets") or {}).get("items", [])[:4]
    if not items:
        return _stack_blocks(slide, card, colors, fonts, start_y=1.9)

    gap = 0.25
    box_w = (11.5 - gap * (len(items) - 1)) / len(items)
    for i, item in enumerate(items):
        x = 0.9 + i * (box_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5),
                                     Inches(box_w), Inches(2.4))
        box.fill.solid(); box.fill.fore_color.rgb = _rgb(colors.get("card", colors["bg"]))
        box.line.color.rgb = _rgb(colors["accent"]); box.line.width = Pt(1.25)
        add_text_box(slide, str(i + 1), Inches(x + 0.25), Inches(2.7), Inches(0.9), Inches(0.7),
                     font_name=fonts["heading"], font_size=32, bold=True, color=colors["accent"])
        add_text_box(slide, str(item), Inches(x + 0.25), Inches(3.45), Inches(box_w - 0.5), Inches(1.3),
                     font_name=fonts["body"], font_size=13, color=colors["text"])


def _render_two_panel(prs, card, colors, fonts, comparison: bool):
    """Two balanced columns. In comparison mode each side gets a header bar."""
    from pptx.enum.shapes import MSO_SHAPE
    slide = _blank_slide(prs, colors)
    add_text_box(slide, card.get("title", ""), Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.9),
                 font_name=fonts["heading"], font_size=26, bold=True, color=colors["text"])

    blocks = card.get("blocks") or []
    table = _first(blocks, "table")
    bullet_blocks = [b for b in blocks if b["type"] == "bullets"][:2]

    # A 2-3 column table reads naturally as two sides; otherwise use two bullet lists.
    if comparison and table and len(table.get("columns", [])) >= 2:
        cols = table["columns"][:2]
        sides = [[r[i] for r in table["rows"]] for i in range(2)]
        headers = cols
    elif len(bullet_blocks) >= 2:
        sides = [bullet_blocks[0]["items"], bullet_blocks[1]["items"]]
        headers = ["", ""]
    else:
        return _stack_blocks(slide, card, colors, fonts, start_y=1.9)

    panel_w = 5.6
    for i in range(2):
        x = 0.9 + i * (panel_w + 0.3)
        y = 2.0
        if comparison and headers[i]:
            hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                         Inches(panel_w), Inches(0.6))
            hdr.fill.solid(); hdr.fill.fore_color.rgb = _rgb(colors["accent"])
            hdr.line.fill.background()
            add_text_box(slide, str(headers[i]), Inches(x + 0.2), Inches(y + 0.1),
                         Inches(panel_w - 0.4), Inches(0.45),
                         font_name=fonts["heading"], font_size=15, bold=True, color=colors["bg"])
            y += 0.85
        for item in sides[i][:6]:
            add_text_box(slide, f"•  {item}", Inches(x + 0.15), Inches(y),
                         Inches(panel_w - 0.3), Inches(0.55),
                         font_name=fonts["body"], font_size=13, color=colors["text"])
            y += 0.58


def _render_closing(prs, card, colors, fonts):
    """Final slide: centred takeaways."""
    slide = _blank_slide(prs, colors)
    title = card.get("title", "")
    size, height, _ = headline_fit(title, box_width_in=10.0, max_font=40)
    add_text_box(slide, title, Inches(1.6), Inches(1.6), Inches(10.0), Inches(max(height, 1.0)),
                 font_name=fonts["heading"], font_size=size, bold=True, color=colors["text"])
    _accent_rule(slide, colors, left=1.6, top=1.6 + max(height, 1.0) + 0.2)

    y = 1.6 + max(height, 1.0) + 0.7
    for block in (card.get("blocks") or []):
        if y > 6.4:
            break
        if block["type"] == "bullets":
            for item in block["items"][:5]:
                add_text_box(slide, f"•  {item}", Inches(1.7), Inches(y), Inches(9.8), Inches(0.55),
                             font_name=fonts["body"], font_size=16, color=colors["text"])
                y += 0.6
        elif block["type"] in ("text", "callout"):
            add_text_box(slide, block["text"][:300], Inches(1.7), Inches(y), Inches(9.8), Inches(0.9),
                         font_name=fonts["body"], font_size=16, color=colors["text"])
            y += 1.0


def _render_card_slide(prs, card: dict, colors: dict, fonts: dict):
    """Dispatches to the renderer for this card's layout. Layouts that are just
    'title plus content' (title_content, text_image, chart, ...) fall through to
    the original stacked renderer, which already handles images and charts."""
    layout = LEGACY_LAYOUTS.get(card.get("layout"), card.get("layout")) or "title_content"
    special = {
        "cover":      _render_cover,
        "section":    _render_section,
        "quote":      _render_quote,
        "stats":      _render_stats,
        "timeline":   _render_timeline,
        "process":    _render_process,
        "closing":    _render_closing,
    }
    if layout in special:
        return special[layout](prs, card, colors, fonts)
    if layout in ("comparison", "two_column"):
        return _render_two_panel(prs, card, colors, fonts, comparison=(layout == "comparison"))
    return _render_stacked_slide(prs, card, colors, fonts)


def _render_stacked_slide(prs, card: dict, colors: dict, fonts: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, colors["bg"])

    # Accent bar + title
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.45), Inches(0.09), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = _rgb(colors["accent"]); bar.line.fill.background()

    title = card.get("title", "")
    font_size, box_height, _ = headline_fit(title, box_width_in=11.6, max_font=28)
    add_text_box(slide, title, Inches(0.85), Inches(0.35), Inches(11.6), Inches(box_height),
                 font_name=fonts["heading"], font_size=font_size, bold=True, color=colors["text"])

    y = 0.35 + box_height + 0.25

    # text_image (formerly image_right): reserve the right column for the image
    blocks = list(card.get("blocks") or [])
    image_block = next((b for b in blocks if b["type"] == "image"), None)
    content_width_in = 12.0
    layout = LEGACY_LAYOUTS.get(card.get("layout"), card.get("layout"))
    if layout == "text_image" and image_block:
        content_width_in = 7.2
        img_bytes = _download_image(image_block.get("url", ""))
        if img_bytes:
            add_picture_fit(slide, img_bytes, Inches(8.0), Inches(y), Inches(4.7), Inches(4.9))
        blocks = [b for b in blocks if b is not image_block]

    _stack_blocks(slide, {"blocks": blocks}, colors, fonts,
                  start_y=y, content_width_in=content_width_in)


def _stack_blocks(slide, card: dict, colors: dict, fonts: dict,
                  start_y: float = 1.9, content_width_in: float = 12.0):
    """Stacks blocks top-down on an existing slide. Shared by the default
    renderer and used as the fallback when a specialised layout is missing the
    block it needs (e.g. a 'timeline' card with no bullets)."""
    y = start_y
    blocks = list(card.get("blocks") or [])
    body_font = fonts["body"]
    for block in blocks:
        if y > 6.6:
            break  # never overflow the slide
        btype = block["type"]

        if btype == "text":
            lines = max(1, len(block["text"]) // 95 + 1)
            h = 0.32 * lines + 0.1
            add_text_box(slide, block["text"], Inches(0.85), Inches(y),
                         Inches(content_width_in - 0.85), Inches(h),
                         font_name=body_font, font_size=14, color=colors["text"])
            y += h + 0.15

        elif btype == "bullets":
            for item in block["items"]:
                add_text_box(slide, f"•  {item}", Inches(0.95), Inches(y),
                             Inches(content_width_in - 1.0), Inches(0.5),
                             font_name=body_font, font_size=14, color=colors["text"])
                y += 0.52
            y += 0.1

        elif btype == "stats":
            items = block["items"]
            col_w = (content_width_in - 0.85) / max(len(items), 1)
            for i, stat in enumerate(items):
                x = 0.85 + i * col_w
                add_text_box(slide, stat["value"], Inches(x), Inches(y), Inches(col_w - 0.2), Inches(0.7),
                             font_name=fonts["heading"], font_size=30, bold=True, color=colors["accent"])
                add_text_box(slide, stat["label"], Inches(x), Inches(y + 0.7), Inches(col_w - 0.2), Inches(0.35),
                             font_name=body_font, font_size=12, bold=True, color=colors["text"])
                if stat.get("sublabel"):
                    add_text_box(slide, stat["sublabel"], Inches(x), Inches(y + 1.05), Inches(col_w - 0.2), Inches(0.3),
                                 font_name=body_font, font_size=10, color=colors["muted"])
            y += 1.6

        elif btype == "chart":
            chart_spec = block["chart"]
            cd = ChartData()
            cd.categories = chart_spec["labels"]
            for ds in chart_spec["datasets"]:
                cd.add_series(ds["label"], [v if v is not None else 0 for v in ds["data"]])
            xl_map = {"bar": XL_CHART_TYPE.BAR_CLUSTERED, "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                      "line": XL_CHART_TYPE.LINE, "pie": XL_CHART_TYPE.PIE}
            h = min(3.4, 6.9 - y)
            if h >= 2.0:
                slide.shapes.add_chart(xl_map.get(chart_spec["type"], XL_CHART_TYPE.COLUMN_CLUSTERED),
                                       Inches(0.85), Inches(y), Inches(content_width_in - 0.85), Inches(h), cd)
                y += h + 0.1
            if block.get("insight") and y < 6.6:
                add_text_box(slide, block["insight"], Inches(0.85), Inches(y),
                             Inches(content_width_in - 0.85), Inches(0.4),
                             font_name=body_font, font_size=12, italic=True, color=colors["muted"])
                y += 0.5

        elif btype == "table":
            cols, rows = block["columns"], block["rows"]
            n_rows = min(len(rows) + 1, int((6.9 - y) / 0.38))
            if n_rows >= 2:
                shape = slide.shapes.add_table(n_rows, len(cols), Inches(0.85), Inches(y),
                                               Inches(content_width_in - 0.85), Inches(0.38 * n_rows))
                table = shape.table
                for c, name in enumerate(cols):
                    cell = table.cell(0, c); cell.text = name
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].runs[0].font.bold = True
                for r in range(1, n_rows):
                    for c in range(len(cols)):
                        cell = table.cell(r, c); cell.text = rows[r - 1][c]
                        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                y += 0.38 * n_rows + 0.15

        elif btype == "callout":
            fill = CALLOUT_FILL[block["style"]]; text_c = CALLOUT_TEXT[block["style"]]
            lines = max(1, len(block["text"]) // 90 + 1)
            h = 0.35 * lines + 0.25
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(y),
                                         Inches(content_width_in - 0.85), Inches(h))
            box.fill.solid(); box.fill.fore_color.rgb = _rgb(fill); box.line.fill.background()
            add_text_box(slide, block["text"], Inches(1.05), Inches(y + 0.08),
                         Inches(content_width_in - 1.3), Inches(h - 0.16),
                         font_name=body_font, font_size=13, bold=True, color=text_c)
            y += h + 0.15

        elif btype == "quote":
            lines = max(1, len(block["text"]) // 80 + 1)
            h = 0.4 * lines
            add_text_box(slide, f'“{block["text"]}”', Inches(1.1), Inches(y),
                         Inches(content_width_in - 1.4), Inches(h),
                         font_name=fonts["heading"], font_size=17, italic=True, color=colors["accent"])
            y += h + 0.05
            if block.get("attribution"):
                add_text_box(slide, f'— {block["attribution"]}', Inches(1.1), Inches(y),
                             Inches(content_width_in - 1.4), Inches(0.35),
                             font_name=body_font, font_size=12, color=colors["muted"])
                y += 0.45

        # The caller already removes the image when it owns the right column,
        # so anything reaching here is meant to sit inline.
        elif btype == "image":
            img_bytes = _download_image(block.get("url", ""))
            h = min(3.2, 6.9 - y)
            if img_bytes and h >= 1.5:
                add_picture_fit(slide, img_bytes, Inches(0.85), Inches(y),
                                Inches(content_width_in - 0.85), Inches(h))
                y += h + 0.15


def _rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _render_title_slide(prs, title: str, subtitle: str, colors: dict, fonts: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, colors["accent"] if colors["bg"] == "FFFFFF" else colors["bg"])
    text_color = "FFFFFF" if colors["bg"] == "FFFFFF" else colors["text"]
    font_size, box_height, _ = headline_fit(title, box_width_in=11.0, max_font=40)
    add_text_box(slide, title, Inches(1.1), Inches(2.6), Inches(11.0), Inches(box_height),
                 font_name=fonts["heading"], font_size=font_size, bold=True, color=text_color)
    if subtitle:
        add_text_box(slide, subtitle, Inches(1.1), Inches(2.6 + box_height + 0.2), Inches(11.0), Inches(0.6),
                     font_name=fonts["body"], font_size=16, color=text_color)


@router.post("/export-deck")
def export_deck(request: ExportDeckRequest):
    """Renders the card deck to PPTX (title slide + one slide per card) on R2."""
    try:
        if not request.cards:
            raise HTTPException(status_code=400, detail="cards cannot be empty.")
        theme  = THEMES.get(request.theme or "aurora_light", THEMES["aurora_light"])
        colors = _theme_colors(request.theme or "aurora_light")
        fonts  = theme["fonts"]

        prs = Presentation()
        prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

        _render_title_slide(prs, request.title, request.subtitle or "", colors, fonts)
        for raw_card in request.cards:
            # enable_images=False prevents new Unsplash fetches; image blocks
            # with an already-resolved url are kept by normalize_card
            card = normalize_card(raw_card, enable_images=False)
            _render_card_slide(prs, card, colors, fonts)

        buf = io.BytesIO(); prs.save(buf)
        pptx_bytes = buf.getvalue()

        safe_title = re.sub(r"[^A-Za-z0-9_-]+", "_", request.title)[:50] or "deck"
        key = f"presentations/{uuid.uuid4()}/{safe_title}.pptx"
        url = upload_to_r2(pptx_bytes, key,
                           "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        print(f"[export-deck] {len(request.cards)} cards -> {key} ({len(pptx_bytes)/1024:.1f}KB)")
        return {"success": True, "pptx_url": url,
                "slide_count": len(request.cards) + 1,
                "size_kb": round(len(pptx_bytes) / 1024, 1)}
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"EXPORT-DECK ERROR: {e}"); print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Deck export failed: {e}")

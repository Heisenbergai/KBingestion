import os
import re
import io
import csv
import uuid
import fitz
import docx
import pptx
import openpyxl
import threading
import time
import ai
from datetime import datetime, timezone
from fastapi import APIRouter, HTTPException, Depends
from auth import AuthContext, current_user
from pydantic import BaseModel
from typing import Optional
from supabase import create_client
from langchain_text_splitters import RecursiveCharacterTextSplitter
from dotenv import load_dotenv
import httpx

load_dotenv()

router = APIRouter()

supabase = create_client(
    os.getenv("SUPABASE_URL"),
    os.getenv("SUPABASE_SERVICE_KEY")
)


# ── In-memory ingestion job store ───────────────────────────────────────────────
# Same pattern as bg_videos.py: Railway is stateless, jobs live in memory only.
# Lovable polls GET /ingest-status/{job_id} and owns persistence of the final
# "processed" flag on the document row. If Railway restarts mid-job, the job
# disappears — Lovable should treat a 404 on a job it was polling as "failed,
# please re-upload".
INGEST_JOBS: dict[str, dict] = {}
MAX_JOBS_KEPT = 200


def _prune_jobs():
    """Keeps the in-memory job store bounded. Oldest finished jobs go first."""
    if len(INGEST_JOBS) <= MAX_JOBS_KEPT:
        return
    finished = [
        (jid, j) for jid, j in INGEST_JOBS.items()
        if j["status"] in ("completed", "failed")
    ]
    finished.sort(key=lambda item: item[1].get("finished_at") or "")
    for jid, _ in finished[: len(INGEST_JOBS) - MAX_JOBS_KEPT]:
        INGEST_JOBS.pop(jid, None)


class IngestRequest(BaseModel):
    document_id:  str
    signed_url:   str
    mime_type:    str
    file_name:    str
    asset_id:     str
    workspace_id: str   # ← REQUIRED — isolates chunks per company
    # Provenance for hybrid-search tiering (Phase 1 / company brain).
    # Uploaded documents are tier-1 (most trusted); connectors override
    # these (e.g. Slack notes = tier 3, meeting notes = tier 2).
    source_type:  Optional[str] = "document"
    source_tier:  Optional[int] = 1
    doc_date:     Optional[str] = None   # ISO date; falls back to created_at
    # Document classification (Phase C) — mirrored onto every chunk so a later
    # phase can filter retrieval by them. Defaults match knowledge_items' own
    # column defaults, in case an older frontend build omits these fields.
    sensitivity:      Optional[str] = "internal"
    authority:        Optional[str] = "working"
    doc_class:        Optional[str] = None
    lifecycle_status: Optional[str] = "active"
    # H-0: optional, default None — an older frontend build that omits them
    # behaves exactly as before (these were null on every chunk until now).
    effective_from:   Optional[str] = None
    valid_until:      Optional[str] = None
    superseded_by:    Optional[str] = None
    # Phase D — frontend-resolved name of the department that owns the
    # uploading folder, if any (Railway cannot query knowledge_folders/
    # departments itself, same reason linked_folder_ids has to be resolved
    # frontend-side). Used only as a rules-engine hint for doc_class.
    folder_department_name: Optional[str] = None
    # wait=True runs synchronously and returns the full result in one response
    # (old behavior — fine for small files and curl testing). Default is
    # background mode: returns a job_id immediately so large documents can't
    # hit the HTTP timeout, and Lovable polls /ingest-status/{job_id}.
    wait: Optional[bool] = False


def download_file(signed_url: str) -> bytes:
    response = httpx.get(signed_url, timeout=120, follow_redirects=True)
    response.raise_for_status()
    return response.content


# ── Extraction helpers ──────────────────────────────────────────────────────────

def _table_to_lines(rows: list[list[str]]) -> list[str]:
    """
    Renders a table as readable lines. If the table has a header row,
    each data row becomes 'Header: Value | Header: Value' — the same
    readable format as extract_xlsx, because that embeds/retrieves better
    than raw pipe-separated cells.
    """
    rows = [[("" if c is None else str(c).strip()) for c in row] for row in rows]
    rows = [r for r in rows if any(r)]
    if not rows:
        return []
    if len(rows) == 1:
        return [" | ".join(c for c in rows[0] if c)]
    header = rows[0]
    lines = []
    for row in rows[1:]:
        pairs = [f"{h}: {v}" for h, v in zip(header, row) if v]
        if pairs:
            lines.append(" | ".join(pairs))
    return lines


def extract_pdf(file_bytes: bytes) -> str:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = "\n".join(page.get_text() for page in doc)
    if doc.page_count > 0 and len(text.strip()) < 30:
        raise ValueError(
            "This PDF appears to be scanned/image-only (no text layer). "
            "OCR is not supported yet — please upload a text-based PDF "
            "or the original document (DOCX/PPTX)."
        )
    return text


def extract_docx(file_bytes: bytes) -> str:
    """
    Extracts paragraphs AND tables, in document order where possible.
    Plain-paragraph extraction silently drops tables, which is where
    policy docs and process docs keep their most important content.
    """
    document = docx.Document(io.BytesIO(file_bytes))
    parts = []
    try:
        # python-docx >= 1.1 yields paragraphs and tables interleaved in order
        for block in document.iter_inner_content():
            if isinstance(block, docx.table.Table):
                rows = [[cell.text for cell in row.cells] for row in block.rows]
                parts.extend(_table_to_lines(rows))
            elif block.text.strip():
                parts.append(block.text)
    except AttributeError:
        # Older python-docx — paragraphs first, then all tables
        parts = [p.text for p in document.paragraphs if p.text.strip()]
        for table in document.tables:
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            parts.extend(_table_to_lines(rows))
    return "\n".join(parts)


def _pptx_shape_texts(shape) -> list[str]:
    """Recursively pulls text out of a shape: text frames, tables, groups."""
    texts = []
    # Grouped shapes contain child shapes
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            texts.extend(_pptx_shape_texts(child))
        return texts
    if getattr(shape, "has_table", False):
        rows = [
            [cell.text for cell in row.cells]
            for row in shape.table.rows
        ]
        texts.extend(_table_to_lines(rows))
        return texts
    if hasattr(shape, "text") and shape.text.strip():
        texts.append(shape.text)
    return texts


def extract_pptx(file_bytes: bytes) -> str:
    """
    Extracts per slide: all shape text (including inside groups), table
    contents, and speaker notes. Notes often carry the actual explanation
    of a slide, so they matter a lot for search quality.
    """
    presentation = pptx.Presentation(io.BytesIO(file_bytes))
    slides_text = []
    for i, slide in enumerate(presentation.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            slide_content.extend(_pptx_shape_texts(shape))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_content.append(f"(Speaker notes: {notes})")
        if slide_content:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(slide_content))
    return "\n\n".join(slides_text)


_BANNER_SCAN_LIMIT = 3


def _find_header_row(rows: list) -> Optional[int]:
    """
    Index of the row that holds the real column names, or None.

    Default is the first non-empty row — the rule both xlsx extractors have
    always used, and the one they must keep agreeing on so the prose view and
    the structured view never disagree about where data starts.

    BANNER ROWS (P1-14). Real business spreadsheets routinely open with a
    merged title cell — "MASTER BUDGET SUMMARY" — above the actual header.
    Taking that as the header produced ["MASTER BUDGET SUMMARY","col_1",
    "col_2",...] and pushed the true column names down into the first DATA
    row. Confirmed on 100% of a real 5-workbook corpus, every sheet. It broke
    the metric card (a picker offering col_2/col_3 is unusable) and weakened
    spreadsheet retrieval, since every numeric row embedded with meaningless
    labels attached.

    So: skip leading rows that have exactly ONE non-empty cell, but ONLY if a
    row with 2+ non-empty cells appears within the next few rows. A genuinely
    single-column sheet never finds one and keeps its original header — the
    fix cannot eat a real header. Bounded scan so a long single-column
    preamble can't run away.

    BUG FOUND AND FIXED 2026-08-13, before this ever ran against real data.
    The original loop treated "the very next non-empty row is also narrow"
    as proof nothing wider follows AT ALL, and gave up immediately — so a
    two-row banner (a merged title, then a narrower subtitle/date line, THEN
    the real header) was left unresolved even though _BANNER_SCAN_LIMIT was
    clearly meant to allow skipping more than one such row. Caught by a
    stress test built from the real corpus's "Executive Dashboard" sheet
    shape (title + "As of Q2" line + header), not assumed safe from the
    original single-banner-row fixture alone. Fixed by collecting up to
    _BANNER_SCAN_LIMIT consecutive narrow rows as banner candidates and only
    giving up once that budget is exhausted or the sheet genuinely ends,
    rather than bailing on the first narrow row encountered mid-scan.
    """
    first_non_empty = None
    for i, row in enumerate(rows):
        if row and any(c is not None and str(c).strip() != "" for c in row):
            first_non_empty = i
            break
    if first_non_empty is None:
        return None

    def width(row) -> int:
        return sum(1 for c in row if c is not None and str(c).strip() != "")

    idx = first_non_empty
    banner_rows_seen = 0
    while width(rows[idx]) == 1 and banner_rows_seen < _BANNER_SCAN_LIMIT:
        banner_rows_seen += 1
        nxt = None
        for j in range(idx + 1, len(rows)):
            if rows[j] and any(c is not None and str(c).strip() != "" for c in rows[j]):
                nxt = j
                break
        if nxt is None:
            # Nothing else in the sheet at all -> keep the original header,
            # never eat the only row that exists.
            return first_non_empty
        idx = nxt

    if width(rows[idx]) >= 2:
        return idx
    # Exhausted the scan budget (or hit another narrow row right at the
    # boundary) without ever reaching a wide row -> this sheet's "header"
    # really is just narrow, or the banner runs deeper than we're willing to
    # guess through. Keep the original first row, same fail-safe as before.
    return first_non_empty


def extract_xlsx(file_bytes: bytes) -> str:
    """
    Reads every sheet in the workbook and converts each row to a readable
    'Column: Value' line, grouped by sheet. Skips fully empty rows.
    This is intentionally verbose/readable rather than compact CSV, since
    the text goes straight into embedding + LLM context — readable text
    embeds and retrieves better than raw comma-separated values.
    """
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
    sections = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        header_idx = _find_header_row(rows)
        if header_idx is None:
            continue
        header = [str(c).strip() if c is not None else f"col{i}"
                  for i, c in enumerate(rows[header_idx])]
        data_rows = [r for r in rows[header_idx + 1:]
                     if not all(cell is None for cell in r)]

        if not data_rows:
            continue

        lines = [f"[Sheet: {sheet_name}]"]
        for row in data_rows:
            pairs = []
            for col_name, value in zip(header, row):
                if value is None or str(value).strip() == "":
                    continue
                pairs.append(f"{col_name}: {value}")
            if pairs:
                lines.append(" | ".join(pairs))

        if len(lines) > 1:
            sections.append("\n".join(lines))

    wb.close()
    return "\n\n".join(sections)


def extract_csv(file_bytes: bytes) -> str:
    """Same 'Column: Value' readable format as extract_xlsx, for consistency."""
    text = file_bytes.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return ""

    header = rows[0]
    lines = []
    for row in rows[1:]:
        pairs = []
        for col_name, value in zip(header, row):
            if value is None or str(value).strip() == "":
                continue
            pairs.append(f"{col_name}: {value}")
        if pairs:
            lines.append(" | ".join(pairs))

    return "\n".join(lines)


def extract_xlsx_tables(file_bytes: bytes, max_rows_per_sheet: int = 5000) -> list[dict]:
    """
    The STRUCTURED counterpart to extract_xlsx — same workbook, types kept.

    WHY THIS EXISTS. extract_xlsx flattens every row into prose
    ("Quarter: Q3 | Revenue: 4200000") because that embeds and retrieves better,
    which is correct for RAG. But it means a spreadsheet — the one input that
    arrives already structured and already typed — has its structure thrown away
    at the last step. A cell containing 4200000 IS the number: no LLM reads it,
    so there is nothing to hallucinate. That is the highest-confidence data this
    system can obtain, and it was being discarded.

    DELIBERATELY A SEPARATE FUNCTION. extract_xlsx is on the live path for every
    upload and every Google Drive sync; its signature, return type and exact
    output text are unchanged and must stay that way. This one is additive and
    its failure is never allowed to abort an ingest (see the caller).

    Note on formulas: openpyxl's data_only=True returns the value Excel/Sheets
    last CACHED, not a computed one. Files saved by real Excel or exported from
    Google Sheets carry those cached values; a workbook generated
    programmatically and never opened does not, and those cells read as None.
    That is a property of the format, not a bug here.
    """
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
    tables: list[dict] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        # Shared header rule (_find_header_row) — the SAME function extract_xlsx
        # calls, so the two views of one sheet cannot disagree about where the
        # data starts. It also skips a leading merged title banner; see P1-14 in
        # that helper's docstring.
        header_idx = _find_header_row(rows)
        if header_idx is None:
            continue
        header = [str(c).strip() if c is not None else f"col_{j}"
                  for j, c in enumerate(rows[header_idx])]

        records: list[dict] = []
        for row in rows[header_idx + 1:]:
            if not row or all(c is None or str(c).strip() == "" for c in row):
                continue  # blank row, same skip rule as extract_xlsx
            record = {}
            for col_name, value in zip(header, row):
                if value is None:
                    continue
                # Keep native types. Dates become ISO strings only because the
                # destination is JSONB, which has no date type.
                if hasattr(value, "isoformat"):
                    record[col_name] = value.isoformat()
                else:
                    record[col_name] = value
            if record:
                records.append(record)
            if len(records) >= max_rows_per_sheet:
                break

        if not records:
            continue

        # Which columns are usable as metrics — the whole point of keeping types.
        numeric_columns = [
            col for col in header
            if any(isinstance(r.get(col), (int, float)) and not isinstance(r.get(col), bool)
                   for r in records)
        ]

        tables.append({
            "sheet":           sheet_name,
            "headers":         header,
            "rows":            records,
            "row_count":       len(records),
            "numeric_columns": numeric_columns,
        })

    return tables


def is_spreadsheet(mime_type: str, file_name: str) -> bool:
    """Mirrors extract_text's own xlsx branch so the two cannot drift apart."""
    return (
        mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        or (file_name or "").lower().endswith((".xlsx", ".xlsm"))
    )


def store_document_tables(tables: list[dict], document_id: str, workspace_id: str,
                          sensitivity: str = "internal") -> int:
    """
    Persists the structured view. Re-ingesting a document REPLACES its tables
    rather than accumulating copies — a re-processed spreadsheet should not end
    up with two conflicting versions of the same sheet, which is exactly the
    kind of ambiguity this whole feature exists to remove.

    Returns how many sheets were stored. Raises on failure; the caller decides
    whether that is fatal (it is not — see process_document_bytes).
    """
    if not tables:
        return 0

    supabase.table("document_tables").delete().eq("document_id", document_id).execute()

    payload = [{
        "document_id":     document_id,
        "workspace_id":    workspace_id,
        "sheet_name":      t["sheet"],
        "headers":         t["headers"],
        "rows":            t["rows"],
        "row_count":       t["row_count"],
        "numeric_columns": t["numeric_columns"],
        # Mirrored from the document so the metrics endpoint can filter by the
        # caller's sensitivity ladder server-side — the vector DB cannot see
        # knowledge_items' RLS. Same mirroring Phase C did for document_chunks.
        "sensitivity":     sensitivity or "internal",
    } for t in tables]

    supabase.table("document_tables").insert(payload).execute()
    print(f"[ingest] stored {len(payload)} structured sheet(s) for document {document_id}")
    return len(payload)


def extract_text(file_bytes: bytes, mime_type: str, file_name: str) -> str:
    name_lower = file_name.lower()

    if mime_type == "application/pdf" or name_lower.endswith(".pdf"):
        return extract_pdf(file_bytes)

    elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" \
         or name_lower.endswith(".docx"):
        return extract_docx(file_bytes)

    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" \
         or name_lower.endswith(".pptx"):
        return extract_pptx(file_bytes)

    elif mime_type in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ) or name_lower.endswith((".xlsx", ".xlsm")):
        return extract_xlsx(file_bytes)

    elif mime_type == "text/csv" or name_lower.endswith(".csv"):
        return extract_csv(file_bytes)

    elif mime_type == "text/plain" or name_lower.endswith(".txt"):
        return file_bytes.decode("utf-8", errors="ignore")

    else:
        raise ValueError(
            f"Unsupported file type: {mime_type or 'unknown'} ({file_name}). "
            f"Supported: PDF, DOCX, PPTX, XLSX, CSV, TXT."
        )


def extract_doc_date(file_bytes: bytes, mime_type: str, file_name: str) -> Optional[str]:
    """
    Best-effort extraction of a document's real authored/creation date from
    its OWN embedded file metadata -- never invented, never derived from the
    filename or upload time (that's exactly what doc_date=None already
    correctly falls back to via `coalesce(doc_date, created_at)` in
    match_chunks_hybrid's freshness scoring).

    Phase 1 Step 6, forward-only: this only ever returns a genuinely embedded
    date or None. A missing property or a parse failure is exactly
    equivalent to the file never having supplied a doc_date at all --
    retrieval's existing created_at fallback already handles that case
    correctly, so failing here costs nothing beyond not-yet-having-a-date,
    same as before this function existed.
    """
    name_lower = file_name.lower()
    try:
        if mime_type == "application/pdf" or name_lower.endswith(".pdf"):
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            raw = (doc.metadata or {}).get("creationDate") or ""
            # PDF date format: D:YYYYMMDDHHmmSS(+-HH'mm') -- only the digit
            # portion is used; the timezone-offset suffix is ignored rather
            # than parsed, since malformed offsets are common in the wild
            # and getting this wrong would be worse than a UTC approximation.
            m = re.match(r"D:(\d{4})(\d{2})(\d{2})(\d{2})?(\d{2})?(\d{2})?", raw)
            if not m:
                return None
            y, mo, d, h, mi, s = (m.group(i) for i in range(1, 7))
            return datetime(
                int(y), int(mo), int(d), int(h or 0), int(mi or 0), int(s or 0),
                tzinfo=timezone.utc,
            ).isoformat()

        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" \
             or name_lower.endswith(".docx"):
            created = docx.Document(io.BytesIO(file_bytes)).core_properties.created
            return created.isoformat() if created else None

        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" \
             or name_lower.endswith(".pptx"):
            created = pptx.Presentation(io.BytesIO(file_bytes)).core_properties.created
            return created.isoformat() if created else None

        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ) or name_lower.endswith((".xlsx", ".xlsm")):
            created = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True).properties.created
            return created.isoformat() if created else None

    except Exception as e:
        print(f"[ingest] doc_date extraction skipped for {file_name}: {e}")
    return None


def clean_text(text: str) -> str:
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r' {2,}', ' ', text)
    text = re.sub(r'\t', ' ', text)
    return text.strip()


def chunk_text(text: str) -> list[str]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    return [c.page_content for c in splitter.create_documents([text])]


# Structural markers the extractors themselves emit. These are the most
# reliable section signals available, because we wrote them: extract_pptx
# emits "[Slide 3]", extract_xlsx "[Sheet: Research Budget]", and
# extract_pdf pages similarly. Matched before any heuristic.
_SECTION_MARKER = re.compile(r'^\s*\[(Slide \d+|Sheet: [^\]]+|Page \d+)\]\s*$', re.M)

# A heading-looking line: short, not sentence-punctuated, not a bullet, and
# either Title Case, ALL CAPS, or numbered ("3. Executive Summary").
_HEADING_LINE = re.compile(
    r'^\s*(?:\d+[.)]\s+)?([A-Z][^\n]{2,79})\s*$'
)


# Words that are legitimately lowercase inside a Title Case heading, so they
# must not drag its score down ("Statement of Work", "Return to Office").
_MINOR_WORDS = {"a", "an", "and", "as", "at", "by", "for", "from", "in", "of",
                "on", "or", "the", "to", "vs", "with"}


def _looks_like_heading(line: str) -> bool:
    s = line.strip()
    if not s or len(s) > 80:
        return False
    if s[-1] in ".,;:?!":          # a sentence, not a heading
        return False
    if s[0] in "-*•–":             # a bullet
        return False
    if not _HEADING_LINE.match(s):
        return False

    letters = [c for c in s if c.isalpha()]
    if not letters:
        return False

    # ALL CAPS banner ("REVENUE SUMMARY").
    if sum(1 for c in letters if c.isupper()) / len(letters) >= 0.60:
        return True

    # Title Case ("Executive Summary", "3. Platform Architecture"). Judged on
    # WORD-INITIAL capitals, not the letter ratio: "Executive Summary" is only
    # 12% uppercase by letter and would be wrongly rejected by a letter-ratio
    # rule -- which is exactly what the tests caught here.
    body = re.sub(r'^\s*\d+[.)]\s+', '', s)          # drop a "3." prefix
    words = [w for w in re.split(r'\s+', body) if any(c.isalpha() for c in w)]
    significant = [w for w in words if w.lower().strip(".,&|") not in _MINOR_WORDS]
    if not significant:
        return False
    capitalised = sum(1 for w in significant if w[0].isupper())
    return capitalised / len(significant) >= 0.60


def extract_section_label(chunk: str, previous: Optional[str] = None) -> Optional[str]:
    """
    Best-effort section/heading label for ONE chunk.

    WHY THIS IS DERIVED PER-CHUNK RATHER THAN BY RE-CHUNKING ON HEADINGS.
    Changing chunk boundaries would change every embedding in the corpus and
    force a full re-ingest of every document — an enormous, irreversible
    regression risk for a metadata improvement. This reads the chunk the
    splitter already produced and labels it, so embeddings, chunk_index and
    retrieval behaviour are all bit-for-bit unchanged. Purely additive.

    `previous` carries the last known section forward, so a chunk that begins
    mid-section (very common with 100-char overlap) inherits its heading
    instead of reporting None.

    Returns None when nothing convincing is found — an honest null is better
    than a confidently wrong heading, which would mislead both the UI and any
    future "related topics" hop.
    """
    if not chunk:
        return previous

    marker = _SECTION_MARKER.search(chunk)
    if marker:
        label = marker.group(1).strip()
        # For a slide/page, the line right after the marker is usually its
        # real title, which is far more useful than "Slide 7".
        tail = chunk[marker.end():].lstrip("\n")
        first_line = tail.split("\n", 1)[0].strip() if tail else ""
        if first_line and _looks_like_heading(first_line):
            return f"{label} — {first_line}"[:200]
        return label[:200]

    for line in chunk.split("\n")[:6]:   # headings sit at the top of a chunk
        if _looks_like_heading(line):
            return line.strip()[:200]

    return previous


def embed_chunks(chunks: list[str], on_progress=None,
                 workspace_id: Optional[str] = None, feature: str = "document_ingestion",
                 ) -> list[list[float]]:
    """
    Embeds via AWS Bedrock Titan v2 (see ai.py). Titan takes one text per
    API call, so we process in batches of 25 purely for progress
    reporting. Throttling retries are handled by boto3's adaptive retry
    mode inside ai.embed_texts. Any hard failure raises — partial/corrupt
    documents never get stored.

    on_progress(embedded_count) is called after each batch so the job
    status endpoint can report live progress to the frontend.

    workspace_id/feature: token-usage attribution. `feature` defaults to
    'document_ingestion' (the /ingest caller); brain_connectors.py's
    filtration pipeline passes 'filtration' instead, since embedding a
    distilled note is a different cost driver worth telling apart on the
    dashboards.
    """
    all_embeddings = []
    batch_size = 25

    for i in range(0, len(chunks), batch_size):
        batch = chunks[i:i + batch_size]
        try:
            all_embeddings.extend(
                ai.embed_texts(batch, workspace_id=workspace_id, feature=feature)
            )
        except Exception as e:
            print(f"[ingest] Embedding error on batch {i}-{i+batch_size}: {e}")
            raise RuntimeError(
                f"Bedrock embedding failed on batch {i}-{i+batch_size}: {e}"
            )
        if on_progress:
            on_progress(len(all_embeddings))
    return all_embeddings


# ── Document classification (Phase D) ────────────────────────────────────────────
# Rules engine runs first (free, deterministic); the LLM (ai.chat_json, same
# convention as brain_connectors.py's classify_batch) fills in whatever rules
# didn't confidently resolve, and always runs for sensitivity specifically —
# that's the security-relevant axis and deserves a real read of the content,
# not just a filename/folder guess. Fail-safe throughout: any failure here
# falls back to this project's own defaults rather than aborting ingestion,
# matching how every other "LLM judgment call" in this codebase behaves.

VALID_SENSITIVITY = {"public", "internal", "confidential", "restricted"}
VALID_AUTHORITY   = {"canonical", "official", "working", "reference", "informal"}
VALID_CLASS       = {"financial", "strategy", "policy_sop", "legal", "product",
                      "people", "sales_marketing", "research_reference", "meeting"}
VALID_LIFECYCLE   = {"draft", "active", "under_review", "superseded", "archived"}

CLASS_FILENAME_HINTS: list[tuple[tuple[str, ...], str]] = [
    (("invoice", "budget", "revenue", "expense", "financial"), "financial"),
    (("policy", "sop", "procedure"), "policy_sop"),
    (("contract", "nda", "agreement", "legal"), "legal"),
    (("roadmap", "strategy", "okr", "vision"), "strategy"),
    (("resume", "cv", "offer letter", "job description"), "people"),
    (("pitch", "proposal", "campaign", "brochure"), "sales_marketing"),
    (("meeting", "notes", "minutes", "standup"), "meeting"),
]

# Best-effort — the project's 8 default department names, not exhaustive.
DEPARTMENT_TO_CLASS = {
    "finance": "financial", "sales": "sales_marketing", "marketing": "sales_marketing",
    "legal": "legal", "people": "people", "engineering": "product", "product": "product",
    "leadership": "strategy",
}

BEDROCK_CLASSIFY_MODEL = os.getenv("BEDROCK_CLASSIFY_MODEL") or None  # falls back to ai.CHAT_MODEL

CLASSIFY_SYSTEM = """You are a document classification assistant for a company's internal knowledge base.
Given a document's title and content, classify it along four independent axes.

Sensitivity - who should be able to see it:
  public        - fine for anyone, including outside the company
  internal      - fine for all employees (default if unsure)
  confidential  - sensitive; only managers/admins should see it (financial detail, internal strategy, PII)
  restricted    - highly sensitive (legal disputes, executive compensation, security incidents, M&A); owner-only

Authority - how much to trust it as current truth:
  canonical  - the definitive, official source of truth
  official   - approved and published, but not THE canonical reference
  working    - a draft or working document, not yet finalized
  reference  - background/reference material, not a source of truth
  informal   - casual notes, not vetted

Class - what kind of document it is: financial, strategy, policy_sop, legal, product, people,
  sales_marketing, research_reference, meeting

Lifecycle - is it still current: draft, active, under_review, superseded, archived

Respond ONLY with valid JSON, no markdown fences:
{"sensitivity": "...", "authority": "...", "doc_class": "..." or null, "lifecycle_status": "...", "confidence": "high"|"medium"|"low"}"""


def _rules_engine_classify(
    title: str, source_type: str, department_hint: Optional[str],
) -> tuple[dict, list[str]]:
    """Free, deterministic first pass. Only ever resolves doc_class today —
    the axis with genuinely deterministic signals (source_type, filename
    keywords, folder department). Returns (partial_result, signals)."""
    result: dict = {}
    signals: list[str] = []

    if source_type in ("meeting", "slack", "note"):
        result["doc_class"] = "meeting"
        signals.append(f"source_type:{source_type}")

    if "doc_class" not in result:
        lower_title = title.lower()
        for keywords, doc_class in CLASS_FILENAME_HINTS:
            if any(k in lower_title for k in keywords):
                result["doc_class"] = doc_class
                signals.append(f"filename_match:{keywords[0]}")
                break

    if "doc_class" not in result and department_hint:
        mapped = DEPARTMENT_TO_CLASS.get(department_hint.strip().lower())
        if mapped:
            result["doc_class"] = mapped
            signals.append(f"department_hint:{department_hint}")

    return result, signals


def classify_document(
    title: str, raw_text: str, source_type: str,
    department_hint: Optional[str] = None, workspace_id: Optional[str] = None,
) -> dict:
    """
    Returns {"sensitivity", "authority", "doc_class", "lifecycle_status",
    "confidence", "signals"}. Never raises — a classification failure falls
    back to this project's own defaults (internal/working/None/active, low
    confidence) rather than aborting ingestion, same as every other
    "LLM judgment call" site in this codebase (classify_batch, etc).
    """
    rules_result, signals = _rules_engine_classify(title, source_type, department_hint)
    defaults = {
        "sensitivity": "internal", "authority": "working",
        "doc_class": None, "lifecycle_status": "active", "confidence": "low",
    }

    if not raw_text.strip():
        # Nothing to read (e.g. an image) - trust rules alone, nothing else to go on.
        return {**defaults, **rules_result, "confidence": "high" if rules_result else "low", "signals": signals}

    try:
        verdict = ai.chat_json(
            messages=[{"role": "user", "content":
                       f"Title: {title}\n\nContent (may be truncated):\n{raw_text[:4000]}"}],
            system=CLASSIFY_SYSTEM, max_tokens=300, temperature=0.2,
            model=BEDROCK_CLASSIFY_MODEL, workspace_id=workspace_id, feature="classification",
        )
    except Exception as e:
        print(f"[classify] LLM classification failed (falling back to rules/defaults): {e}")
        signals.append("llm_failed")
        return {**defaults, **rules_result, "signals": signals}

    if not isinstance(verdict, dict):
        signals.append("llm_invalid_response")
        return {**defaults, **rules_result, "signals": signals}

    result = dict(defaults)
    if verdict.get("sensitivity") in VALID_SENSITIVITY:
        result["sensitivity"] = verdict["sensitivity"]
    if verdict.get("authority") in VALID_AUTHORITY:
        result["authority"] = verdict["authority"]
    if verdict.get("lifecycle_status") in VALID_LIFECYCLE:
        result["lifecycle_status"] = verdict["lifecycle_status"]
    if verdict.get("confidence") in ("high", "medium", "low"):
        result["confidence"] = verdict["confidence"]
    # Rules-engine doc_class wins over the LLM's guess - a deterministic
    # keyword/source_type match is more trustworthy than an LLM read for
    # this one axis specifically.
    if "doc_class" in rules_result:
        result["doc_class"] = rules_result["doc_class"]
    elif verdict.get("doc_class") in VALID_CLASS:
        result["doc_class"] = verdict["doc_class"]

    result["signals"] = signals if signals else ["llm_classified"]
    return result


# ── Core pipeline (shared by sync and background modes) ─────────────────────────

def process_document(request: IngestRequest, job: Optional[dict] = None) -> dict:
    """
    Full pipeline: download → extract → clean → chunk → embed → store.
    Updates `job` in place (if given) so /ingest-status shows live progress.
    Raises ValueError for user-fixable problems, other exceptions for real errors.
    """
    def set_stage(stage: str):
        if job is not None:
            job["stage"] = stage

    set_stage("downloading")
    file_bytes = download_file(request.signed_url)

    return process_document_bytes(
        file_bytes, document_id=request.document_id, asset_id=request.asset_id,
        workspace_id=request.workspace_id, mime_type=request.mime_type,
        file_name=request.file_name, source_type=request.source_type or "document",
        source_tier=request.source_tier or 1, doc_date=request.doc_date, job=job,
        sensitivity=request.sensitivity or "internal",
        authority=request.authority or "working",
        doc_class=request.doc_class,
        lifecycle_status=request.lifecycle_status or "active",
        folder_department_name=request.folder_department_name,
        effective_from=request.effective_from,
        valid_until=request.valid_until,
        superseded_by=request.superseded_by,
    )


def process_document_bytes(
    file_bytes: bytes, document_id: str, asset_id: str, workspace_id: str,
    mime_type: str, file_name: str, source_type: str = "document",
    source_tier: int = 1, doc_date: Optional[str] = None, job: Optional[dict] = None,
    sensitivity: str = "internal", authority: str = "working",
    doc_class: Optional[str] = None, lifecycle_status: str = "active",
    folder_department_name: Optional[str] = None,
    # H-0: carried so a RE-INGEST doesn't silently wipe an expiry date the user
    # had set. The other four classification fields were already carried for
    # exactly this reason; these three were left behind when they were added,
    # so "Re-run AI Processing" on a document with a valid_until would have
    # dropped it from the chunks while knowledge_items kept it — a silent
    # divergence between the two, which is what the mirroring exists to prevent.
    effective_from: Optional[str] = None, valid_until: Optional[str] = None,
    superseded_by: Optional[str] = None,
) -> dict:
    """
    The extract → clean → chunk → embed → store tail of process_document(),
    factored out so a caller that already HAS the bytes (no signed_url to
    download from) can skip straight to it. First real caller: connector_google
    fetches/exports a Drive file's bytes directly from Google's API — routing
    that through a signed Supabase Storage URL just to satisfy the download
    step would be a pointless extra upload/download round trip.
    """
    def set_stage(stage: str):
        if job is not None:
            job["stage"] = stage

    set_stage("extracting")
    raw_text = extract_text(file_bytes, mime_type, file_name)

    # Phase 1 Step 6: fill doc_date from the file's own embedded metadata
    # when the caller didn't already supply one -- forward-only, never
    # overrides an explicit caller value, never invents a date, never
    # touches already-ingested rows. See extract_doc_date's docstring.
    if doc_date is None:
        doc_date = extract_doc_date(file_bytes, mime_type, file_name)

    # Phase H: a spreadsheet also keeps its STRUCTURE, not just its prose.
    # Best-effort by design — a structured-parse failure must never cost the
    # user their upload, since the text path above has already succeeded and is
    # what powers search. Same fail-safe convention as ai._log_usage().
    if is_spreadsheet(mime_type, file_name):
        try:
            store_document_tables(
                extract_xlsx_tables(file_bytes), document_id, workspace_id,
                sensitivity=sensitivity,
            )
        except Exception as e:
            print(f"[ingest] structured spreadsheet parse skipped for {file_name}: {e}")

    if not raw_text.strip():
        raise ValueError(
            "No text could be extracted. The file may be empty, image-only "
            "(scanned PDF with no OCR layer), or a spreadsheet with no data rows."
        )

    set_stage("classifying")
    classification = classify_document(
        title=file_name, raw_text=raw_text, source_type=source_type,
        department_hint=folder_department_name, workspace_id=workspace_id,
    )
    classification_fields = {
        "proposed_sensitivity":      classification["sensitivity"],
        "proposed_authority":        classification["authority"],
        "proposed_doc_class":        classification["doc_class"],
        "proposed_lifecycle_status": classification["lifecycle_status"],
        "classification_confidence": classification["confidence"],
        "classification_signals":    classification["signals"],
    }
    if job is not None:
        job.update(classification_fields)

    cleaned = clean_text(raw_text)
    chunks  = chunk_text(cleaned)

    if not chunks:
        raise ValueError("Document too short to process.")

    if job is not None:
        job["chunks_total"] = len(chunks)

    print(f"[ingest] {file_name}: {len(chunks)} chunks to embed")

    set_stage("embedding")

    def on_progress(done: int):
        if job is not None:
            job["chunks_embedded"] = done

    embeddings = embed_chunks(chunks, on_progress=on_progress, workspace_id=workspace_id)

    set_stage("storing")

    # Remove any previous chunks for this document first. This makes
    # re-uploading a document safe (no duplicate chunks polluting search)
    # and is what lets the workspace-isolation re-upload step replace old
    # workspace_id=null chunks instead of stacking on top of them.
    supabase.table("document_chunks") \
        .delete() \
        .eq("document_id", document_id) \
        .execute()

    # Section labels, computed in order so each chunk can inherit the last
    # known heading when it starts mid-section. Derived from the chunks the
    # splitter ALREADY produced — no boundary, embedding or index changes.
    section_labels: list[Optional[str]] = []
    _running_section: Optional[str] = None
    for _c in chunks:
        try:
            _running_section = extract_section_label(_c, _running_section)
        except Exception as e:
            # A labelling failure must never cost a user their upload; the
            # chunk simply goes in unlabelled. Same convention as the
            # structured-spreadsheet parse and ai._log_usage.
            print(f"[ingest] section labelling failed for one chunk (non-fatal): {e}")
            _running_section = None
        section_labels.append(_running_section)

    rows = [
        {
            "document_id":  document_id,
            "asset_id":     asset_id,
            "workspace_id": workspace_id,  # ← stored with every chunk
            "content":      chunks[i],
            "embedding":    embeddings[i],
            "chunk_index":  i,
            "source_type":  source_type,
            "source_tier":  source_tier,
            "doc_date":     doc_date,      # None → DB default (created_at)
            "sensitivity":       sensitivity,
            "authority":         authority,
            "doc_class":         doc_class,
            "lifecycle_status":  lifecycle_status,
            "effective_from":    effective_from,
            "valid_until":       valid_until,
            "superseded_by":     superseded_by,
            "metadata": {
                "file_name":    file_name,
                "chunk_index":  i,
                "total_chunks": len(chunks),
                "workspace_id": workspace_id,
                "source_type":  source_type,
                # Step 2 of the knowledge-structuring work: the heading this
                # chunk sits under. Enables "related headings/topics",
                # heading-aware citations, and a future topic-level hop.
                # Null when nothing convincing was found — an honest null
                # beats a confidently wrong heading.
                "section":      section_labels[i],
            }
        }
        for i in range(len(chunks))
    ]

    # Insert in batches of 200 rows — avoids one giant request for very
    # large documents (e.g. big spreadsheets can produce thousands of chunks)
    INSERT_BATCH = 200
    for i in range(0, len(rows), INSERT_BATCH):
        supabase.table("document_chunks").insert(rows[i:i + INSERT_BATCH]).execute()

    return {
        "success":        True,
        "document_id":    document_id,
        "workspace_id":   workspace_id,
        "chunks_created": len(chunks),
        "message":        f"Processed '{file_name}' into {len(chunks)} chunks.",
        **classification_fields,
    }


def log_processing_outcome(document_id: str, workspace_id: str, status: str,
                           chunks_created: int = 0, error: Optional[str] = None,
                           duration_ms: Optional[int] = None,
                           file_name: Optional[str] = None) -> None:
    """
    H-0: appends one row per ingestion attempt, success or failure.

    Best-effort and fail-safe, same convention as ai._log_usage() and
    signals.log_signal(): a logging failure must NEVER turn a SUCCESSFUL
    ingest into a failed one, and must never mask the REAL error on the
    failure path. Every exception is swallowed here.
    """
    try:
        supabase.table("document_processing_log").insert({
            "document_id":    document_id,
            "workspace_id":   workspace_id,
            "status":         status,
            "chunks_created": chunks_created,
            "error":          (error or None) and error[:2000],
            "duration_ms":    duration_ms,
            "file_name":      file_name,
        }).execute()
    except Exception as e:
        print(f"[ingest] processing-log write failed, non-fatal ({status}): {e}")


def _run_ingest_job(job_id: str, request: IngestRequest):
    job = INGEST_JOBS.get(job_id)
    if job is None:
        return
    started = time.monotonic()
    try:
        result = process_document(request, job=job)
        job.update({
            "status":         "completed",
            "stage":          "completed",
            "chunks_created": result["chunks_created"],
            "finished_at":    datetime.now(timezone.utc).isoformat(),
        })
        print(f"[ingest:job {job_id}] Completed — {result['chunks_created']} chunks")
        log_processing_outcome(
            request.document_id, request.workspace_id, "completed",
            chunks_created=result["chunks_created"],
            duration_ms=int((time.monotonic() - started) * 1000),
            file_name=request.file_name,
        )
    except Exception as e:
        import traceback
        print(f"[ingest:job {job_id}] FAILED: {e}")
        print(traceback.format_exc())
        job.update({
            "status":      "failed",
            "error":       str(e),
            "finished_at": datetime.now(timezone.utc).isoformat(),
        })
        # Logged AFTER job.update, so a logging problem can never prevent the
        # real error reaching /ingest-status — which is what the user actually
        # sees. The log is a health signal, never the primary error channel.
        log_processing_outcome(
            request.document_id, request.workspace_id, "failed",
            error=str(e),
            duration_ms=int((time.monotonic() - started) * 1000),
            file_name=request.file_name,
        )


# ── Routes ──────────────────────────────────────────────────────────────────────

@router.post("/ingest")
async def ingest_document(request: IngestRequest,
                          auth: AuthContext = Depends(current_user)):
    """
    Processes a document and stores chunks in the vector DB.
    workspace_id is stored with every chunk — this is what isolates
    each company's data from all other companies.

    Default (background) mode returns immediately:
        { success, job_id, status: "processing" }
    then the frontend polls GET /ingest-status/{job_id} until
    status is "completed" or "failed". This is what makes large
    documents work — the old synchronous mode timed out on them.

    Pass "wait": true for the old synchronous behavior (small files, curl tests).

    Supported formats: PDF, DOCX, PPTX, XLSX, CSV, TXT.
    """
    if not request.workspace_id:
        raise HTTPException(
            status_code=400,
            detail="workspace_id is required. Every document must belong to a workspace."
        )

    # Chunks are written with this workspace_id, so an unauthorised caller could
    # otherwise inject documents into someone else's brain, not just read one.
    auth.assert_workspace(request.workspace_id)

    if request.wait:
        try:
            return process_document(request)
        except ValueError as e:
            raise HTTPException(status_code=400, detail=str(e))
        except Exception as e:
            import traceback
            print(f"INGEST ERROR: {str(e)}")
            print(traceback.format_exc())
            raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")

    _prune_jobs()
    job_id = str(uuid.uuid4())
    INGEST_JOBS[job_id] = {
        "job_id":          job_id,
        "document_id":     request.document_id,
        "workspace_id":    request.workspace_id,
        "file_name":       request.file_name,
        "status":          "processing",
        "stage":           "queued",
        "chunks_total":    None,
        "chunks_embedded": 0,
        "chunks_created":  None,
        "error":           None,
        "started_at":      datetime.now(timezone.utc).isoformat(),
        "finished_at":     None,
    }

    # Plain daemon thread, not BackgroundTasks — the pipeline is blocking
    # (sync HTTP, sync Voyage client, time.sleep backoff) and can run for
    # minutes on big documents; a thread keeps the event loop free.
    threading.Thread(
        target=_run_ingest_job, args=(job_id, request), daemon=True
    ).start()

    return {
        "success":     True,
        "job_id":      job_id,
        "status":      "processing",
        "document_id": request.document_id,
        "message":     f"Processing '{request.file_name}' in the background. "
                       f"Poll /ingest-status/{job_id} for progress.",
    }


@router.get("/ingest-status/{job_id}")
async def ingest_status(job_id: str,
                        auth: AuthContext = Depends(current_user)):
    """
    Live status of a background ingestion job.
    status: processing | completed | failed
    stage:  queued | downloading | extracting | embedding | storing | completed
    While embedding, chunks_embedded / chunks_total gives real progress
    for a percentage bar in the UI.

    The job record carries file_name and workspace_id, so it is authorised like
    any other workspace read. A random job UUID is not a credential.
    """
    job = INGEST_JOBS.get(job_id)
    if job is None:
        raise HTTPException(
            status_code=404,
            detail="Unknown job_id. The server may have restarted — re-upload the document."
        )
    auth.assert_workspace(job["workspace_id"])
    return job


class DocumentMetadataUpdate(BaseModel):
    document_id:      str
    workspace_id:     str
    sensitivity:      Optional[str] = None
    authority:        Optional[str] = None
    doc_class:        Optional[str] = None
    lifecycle_status: Optional[str] = None
    # H-0: the three lifecycle DATE fields, mirrored the same way the four
    # above already are. The Library's classification editor already sends
    # these to knowledge_items; without them here, a document's expiry would
    # sit in the app DB where Railway (and therefore the health checker and
    # retrieval ranking) cannot see it.
    effective_from:   Optional[str] = None
    valid_until:      Optional[str] = None
    superseded_by:    Optional[str] = None


@router.post("/document-metadata")
async def update_document_metadata(request: DocumentMetadataUpdate,
                                   auth: AuthContext = Depends(current_user)):
    """
    Syncs a classification change onto an already-ingested document's chunks.
    Without this, re-classifying a document via the Library's edit UI would
    only update knowledge_items in the app DB — the chunks already sitting in
    the vector DB from the original ingestion would keep stale values, which
    defeats "mirrored to chunks" as an ongoing guarantee (Phase E's retrieval
    filtering would then read the wrong tier for anything reclassified after
    upload). Metadata only — no re-embedding, since content hasn't changed.
    A document with no chunks yet (never ingested, or still processing) is a
    no-op: the next real ingestion will carry the current values anyway.
    """
    auth.assert_workspace(request.workspace_id)

    patch = {
        k: v for k, v in {
            "sensitivity":      request.sensitivity,
            "authority":        request.authority,
            "doc_class":        request.doc_class,
            "lifecycle_status": request.lifecycle_status,
            "effective_from":   request.effective_from,
            "valid_until":      request.valid_until,
            "superseded_by":    request.superseded_by,
        }.items() if v is not None
    }
    if not patch:
        return {"success": True, "updated_chunks": 0}

    result = supabase.table("document_chunks") \
        .update(patch) \
        .eq("document_id", request.document_id) \
        .eq("workspace_id", request.workspace_id) \
        .execute()

    # document_tables mirrors sensitivity for the SAME reason document_chunks
    # does, and must be synced here too. Missing this was a real, confirmed
    # exposure: automated classification (Phase D) raises a spreadsheet to
    # `confidential` AFTER ingest has already written its tables at the
    # then-current tier, so the tables kept `internal` while the document and
    # its chunks became `confidential`. /document-tables and
    # /document-table-rows filter on document_tables.sensitivity, so an
    # employee-tier ladder (`public`,`internal`) could read every cell of a
    # confidential spreadsheet through the metric card while knowledge_items
    # RLS correctly hid the document itself. Found live 2026-08-01 on a real
    # R&D budget sheet, on the first real spreadsheet ever uploaded.
    #
    # Only `sensitivity` exists on document_tables -- the other classification
    # axes are not mirrored there -- so the patch is narrowed rather than
    # reused wholesale, which would 400 on unknown columns.
    updated_tables = 0
    if "sensitivity" in patch:
        try:
            t_result = supabase.table("document_tables") \
                .update({"sensitivity": patch["sensitivity"]}) \
                .eq("document_id", request.document_id) \
                .eq("workspace_id", request.workspace_id) \
                .execute()
            updated_tables = len(t_result.data or [])
        except Exception as e:
            # Never let a tables-sync failure fail the chunk sync that already
            # succeeded -- but do surface it, because a silent failure here is
            # exactly how the stale value got there in the first place.
            print(f"DOCUMENT-METADATA document_tables sync FAILED: {e}")

    return {
        "success": True,
        "updated_chunks": len(result.data or []),
        "updated_tables": updated_tables,
    }


class DocumentDeletedRequest(BaseModel):
    document_ids: list[str]
    workspace_id: str


@router.post("/document-deleted")
async def sync_document_deleted(request: DocumentDeletedRequest,
                                auth: AuthContext = Depends(current_user)):
    """
    P0 fix, 2026-08-13. Soft-deleting a document (knowledge_items.deleted_at,
    app DB) was leaving its document_chunks and document_tables rows fully
    intact in the vector DB -- confirmed live: a distinctive chunk from the
    first document ever actually soft-deleted on this system was still fully
    retrievable by match_chunks_hybrid/match_chunks_workspace after "deletion"
    (a bot or AI Search could still surface and cite it), and its spreadsheet
    stayed fully pickable via /document-tables. This is the same mirroring
    gap /document-metadata already exists to close for reclassification --
    the vector DB cannot see knowledge_items' own deleted_at, so this must be
    synced explicitly, the same way sensitivity/authority/etc. already are.

    Called by the frontend right after a successful soft-delete (item or
    folder). Best-effort by design, same convention as /document-metadata's
    document_tables sync: never raises past a logged warning, since the
    knowledge_items row is already correctly marked deleted regardless of
    whether this sync succeeds -- a failure here is a (logged, fixable)
    retrieval-staleness gap, not a reason to make the user's delete action
    itself fail.
    """
    auth.assert_workspace(request.workspace_id)
    if not request.document_ids:
        return {"success": True, "updated_chunks": 0, "updated_tables": 0}

    now = datetime.now(timezone.utc).isoformat()
    updated_chunks = 0
    updated_tables = 0
    try:
        c_result = supabase.table("document_chunks") \
            .update({"deleted_at": now}) \
            .in_("document_id", request.document_ids) \
            .eq("workspace_id", request.workspace_id) \
            .execute()
        updated_chunks = len(c_result.data or [])
    except Exception as e:
        print(f"DOCUMENT-DELETED chunks sync FAILED: {e}")

    try:
        t_result = supabase.table("document_tables") \
            .update({"deleted_at": now}) \
            .in_("document_id", request.document_ids) \
            .eq("workspace_id", request.workspace_id) \
            .execute()
        updated_tables = len(t_result.data or [])
    except Exception as e:
        print(f"DOCUMENT-DELETED tables sync FAILED: {e}")

    return {"success": True, "updated_chunks": updated_chunks, "updated_tables": updated_tables}

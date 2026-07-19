import os
import io
import uuid
import json
import httpx
import boto3
from botocore.config import Config
from fastapi import APIRouter, HTTPException
from pydantic import BaseModel
from typing import Optional
import ai
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from dotenv import load_dotenv

load_dotenv()

router = APIRouter()

UNSPLASH_KEY     = os.getenv("UNSPLASH_ACCESS_KEY", "")
R2_ACCOUNT_ID    = os.getenv("R2_ACCOUNT_ID")
R2_ACCESS_KEY_ID = os.getenv("R2_ACCESS_KEY_ID")
R2_SECRET_KEY    = os.getenv("R2_SECRET_ACCESS_KEY")
R2_BUCKET        = os.getenv("R2_BUCKET_NAME", "knowledge-videos")
R2_PUBLIC_URL    = os.getenv("R2_PUBLIC_URL", "https://pub-222a48f277db4cb6b9fbdda27a672cc5.r2.dev")

r2 = boto3.client(
    service_name="s3",
    endpoint_url=f"https://{R2_ACCOUNT_ID}.r2.cloudflarestorage.com",
    aws_access_key_id=R2_ACCESS_KEY_ID,
    aws_secret_access_key=R2_SECRET_KEY,
    region_name="auto",
    config=Config(signature_version="s3v4", retries={"max_attempts": 5}),
)

# ── Color palettes ─────────────────────────────────────────────────────────────
PALETTES = {
    "midnight_executive": {
        "primary":   "1E2761",
        "secondary": "CADCFC",
        "accent":    "FFFFFF",
        "dark_text": "1A1A2E",
        "light_text": "F5F5F5",
    },
    "ocean_gradient": {
        "primary":   "065A82",
        "secondary": "1C7293",
        "accent":    "FFFFFF",
        "dark_text": "0D2137",
        "light_text": "F0F8FF",
    },
    "charcoal_minimal": {
        "primary":   "36454F",
        "secondary": "F2F2F2",
        "accent":    "212121",
        "dark_text": "1C1C1C",
        "light_text": "FAFAFA",
    },
}


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Request shapes ─────────────────────────────────────────────────────────────
class BrandingSettings(BaseModel):
    palette:       Optional[str]  = "midnight_executive"
    primary_color: Optional[str]  = None   # override hex e.g. "1E2761"
    company_name:  Optional[str]  = ""
    font:          Optional[str]  = "Calibri"   # safe font


class GeneratePresentationRequest(BaseModel):
    query:          str
    chunks:         list[str]       # retrieved chunks from AI Search
    tone:           Optional[str]   = "professional"
    num_slides:     Optional[int]   = 8
    enable_charts:  Optional[bool]  = True
    enable_images:  Optional[bool]  = True
    branding:       Optional[BrandingSettings] = BrandingSettings()


class EditSlideRequest(BaseModel):
    slide:       dict     # the current slide JSON
    instruction: str      # user's edit instruction
    context:     Optional[str] = ""   # surrounding context if needed


class ExportPptxRequest(BaseModel):
    slides:        list[dict]
    title:         str
    branding:      Optional[BrandingSettings] = BrandingSettings()
    enable_images: Optional[bool] = True   # fetch+embed Unsplash images where slides request them


# ── Helpers ────────────────────────────────────────────────────────────────────
def add_text_box(slide, text: str, left, top, width, height,
                 font_name="Calibri", font_size=14, bold=False,
                 color="1A1A2E", align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_to_rgb(color)
    return txBox


def set_slide_background(slide, hex_color: str):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)


def upload_to_r2(data: bytes, key: str, content_type: str) -> str:
    r2.put_object(Bucket=R2_BUCKET, Key=key, Body=data, ContentType=content_type)
    return f"{R2_PUBLIC_URL}/{key}"


# ── Headline auto-fit (FIX for overflow-into-bullets bug) ──────────────────────
def headline_fit(headline: str, box_width_in: float = 12.0,
                  max_font: int = 32, min_font: int = 21) -> tuple[int, float, int]:
    """
    Picks a font size (32 -> 26 -> 21pt) that keeps a headline to at most
    2 lines, and returns a box height sized to how many lines it actually
    wraps to. Callers use the returned height to position whatever comes
    next (bullets, charts, stats) instead of a fixed Y offset — this is
    what fixes long headlines overflowing into the content below them.

    python-pptx has no text-measurement API, so line count is estimated
    from average character width per font size. The estimate is
    deliberately conservative (slightly overestimates width) so lines
    wrap a little early rather than overflow.

    Returns: (font_size_pt, box_height_inches, estimated_line_count)
    """
    for font_size in (max_font, max_font - 6, min_font):
        avg_char_width_in = (font_size * 0.52) / 72.0
        chars_per_line = max(int(box_width_in / avg_char_width_in), 1)
        line_count = max(1, -(-len(headline) // chars_per_line))  # ceil division
        if line_count <= 2 or font_size == min_font:
            line_count = min(line_count, 3)  # hard cap — never plan for more than 3 lines
            box_height = 0.15 + line_count * (font_size * 1.25) / 72.0
            return font_size, round(box_height, 2), line_count
    return min_font, 1.3, 3


def add_picture_fit(slide, image_bytes: bytes, left, top, box_width, box_height):
    """Places an image centered inside a box, preserving aspect ratio (letterboxed)."""
    with PILImage.open(io.BytesIO(image_bytes)) as im:
        img_w, img_h = im.size
    img_ratio = img_w / img_h
    box_ratio = box_width / box_height

    if img_ratio > box_ratio:
        pic_width  = box_width
        pic_height = int(box_width / img_ratio)
    else:
        pic_height = box_height
        pic_width  = int(box_height * img_ratio)

    pic_left = int(left + (box_width - pic_width) / 2)
    pic_top  = int(top + (box_height - pic_height) / 2)

    stream = io.BytesIO(image_bytes)
    slide.shapes.add_picture(stream, pic_left, pic_top, width=pic_width, height=pic_height)


async def fetch_unsplash_image(query: str) -> Optional[bytes]:
    """Fetches one landscape photo for a query. Best-effort — returns None on any failure."""
    if not UNSPLASH_KEY:
        return None
    try:
        async with httpx.AsyncClient(timeout=10) as client:
            res = await client.get(
                "https://api.unsplash.com/search/photos",
                params={"query": query, "per_page": 1, "orientation": "landscape"},
                headers={"Authorization": f"Client-ID {UNSPLASH_KEY}"}
            )
            data = res.json()
            results = data.get("results") or []
            if not results:
                return None

            photo = results[0]
            img_url = photo["urls"]["regular"]
            img_res = await client.get(img_url)

            # Unsplash API guidelines: register a download event when a photo is used.
            # Best-effort — never fail the presentation over this.
            download_location = photo.get("links", {}).get("download_location")
            if download_location:
                try:
                    await client.get(download_location, headers={"Authorization": f"Client-ID {UNSPLASH_KEY}"})
                except Exception:
                    pass

            return img_res.content
    except Exception as e:
        print(f"[presentation] Unsplash fetch failed for '{query}': {e}")
    return None


async def prefetch_images(slides: list[dict], enable_images: bool) -> dict[int, bytes]:
    """
    Pre-fetches images for every slide that has an image_query, keyed by
    slide index. Done up front (before rendering) since Unsplash calls are
    async and python-pptx rendering is synchronous.
    """
    images: dict[int, bytes] = {}
    if not enable_images:
        return images

    for slide_data in slides:
        query = slide_data.get("image_query")
        idx = slide_data.get("index")
        if not query or idx is None:
            continue
        img_bytes = await fetch_unsplash_image(query)
        if img_bytes:
            images[idx] = img_bytes

    return images


# ── Slide rendering ────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def render_title_slide(prs: Presentation, slide_data: dict, palette: dict, font: str, image_bytes: Optional[bytes] = None):
    layout = prs.slide_layouts[6]  # blank
    slide  = prs.slides.add_slide(layout)
    set_slide_background(slide, palette["primary"])

    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")
    company = slide_data.get("company", "")

    add_text_box(slide, title,
                 Inches(1), Inches(2.2), Inches(11), Inches(1.8),
                 font_name=font, font_size=44, bold=True,
                 color=palette["accent"], align=PP_ALIGN.CENTER)

    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(1), Inches(4.1), Inches(11), Inches(0.9),
                     font_name=font, font_size=22,
                     color=palette["secondary"], align=PP_ALIGN.CENTER)

    if company:
        add_text_box(slide, company,
                     Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                     font_name=font, font_size=14,
                     color=palette["secondary"], align=PP_ALIGN.CENTER)


def render_executive_summary(prs, slide_data: dict, palette: dict, font: str, image_bytes: Optional[bytes] = None):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)

    headline = slide_data.get("headline", "Executive Summary")
    bullets  = slide_data.get("bullets", [])[:3]

    font_size, box_height, _ = headline_fit(headline, box_width_in=12.0, max_font=36)
    add_text_box(slide, headline,
                 Inches(0.6), Inches(0.4), Inches(12), Inches(box_height),
                 font_name=font, font_size=font_size, bold=True,
                 color=palette["dark_text"])

    y = Inches(0.4) + Inches(box_height) + Inches(0.3)
    for bullet in bullets:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.6), y + Inches(0.1),
            Inches(0.3), Inches(0.3)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = hex_to_rgb(palette["primary"])
        circle.line.fill.background()

        add_text_box(slide, bullet,
                     Inches(1.1), y, Inches(11.2), Inches(0.9),
                     font_name=font, font_size=16,
                     color=palette["dark_text"])
        y += Inches(1.4)


def render_content_slide(prs, slide_data: dict, palette: dict, font: str, image_bytes: Optional[bytes] = None):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)

    headline = slide_data.get("headline", "")
    bullets  = slide_data.get("bullets", [])

    font_size, box_height, _ = headline_fit(headline, box_width_in=12.0, max_font=32)
    add_text_box(slide, headline,
                 Inches(0.6), Inches(0.4), Inches(12), Inches(box_height),
                 font_name=font, font_size=font_size, bold=True,
                 color=palette["dark_text"])

    y = Inches(0.4) + Inches(box_height) + Inches(0.3)

    # If an image is available, narrow the bullet column and place the
    # image in the freed-up right-hand space.
    bullet_width = Inches(11.2)
    if image_bytes:
        bullet_width = Inches(6.6)
        img_box_left = Inches(7.6)
        img_box_top  = Inches(1.6)
        add_picture_fit(slide, image_bytes, img_box_left, img_box_top, Inches(5.1), Inches(4.8))

    for bullet in bullets[:5]:
        add_text_box(slide, f"• {bullet}",
                     Inches(0.8), y, bullet_width, Inches(0.75),
                     font_name=font, font_size=15,
                     color=palette["dark_text"])
        y += Inches(0.9)


def render_chart_slide(prs, slide_data: dict, palette: dict, font: str, image_bytes: Optional[bytes] = None):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)

    headline   = slide_data.get("headline", "")
    chart_data = slide_data.get("chart", {})
    insight    = slide_data.get("insight", "")

    font_size, box_height, _ = headline_fit(headline, box_width_in=12.0, max_font=30)
    add_text_box(slide, headline,
                 Inches(0.6), Inches(0.3), Inches(12), Inches(box_height),
                 font_name=font, font_size=font_size, bold=True,
                 color=palette["dark_text"])

    chart_top = Inches(0.3) + Inches(box_height) + Inches(0.2)

    cd = ChartData()
    labels   = chart_data.get("labels", ["Q1", "Q2", "Q3", "Q4"])
    datasets = chart_data.get("datasets", [{"label": "Data", "data": [0, 0, 0, 0]}])
    cd.categories = labels

    for ds in datasets:
        data_points = [v if v is not None else 0 for v in ds.get("data", [])]
        cd.add_series(ds.get("label", "Series"), data_points)

    chart_type_str = chart_data.get("type", "bar")
    chart_type_map = {
        "bar":    XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line":   XL_CHART_TYPE.LINE,
        "pie":    XL_CHART_TYPE.PIE,
    }
    xl_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    slide.shapes.add_chart(
        xl_type,
        Inches(0.6), chart_top,
        Inches(8.5), Inches(7.3) - chart_top,
        cd
    )

    if insight:
        add_text_box(slide, f"💡 {insight}",
                     Inches(9.4), chart_top + Inches(1.2), Inches(3.5), Inches(2.5),
                     font_name=font, font_size=14, italic=True,
                     color=palette["primary"])


def render_big_stat_slide(prs, slide_data: dict, palette: dict, font: str, image_bytes: Optional[bytes] = None):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)

    headline = slide_data.get("headline", "")
    stats    = slide_data.get("stats", [])[:3]

    font_size, box_height, _ = headline_fit(headline, box_width_in=12.0, max_font=30)
    add_text_box(slide, headline,
                 Inches(0.6), Inches(0.3), Inches(12), Inches(box_height),
                 font_name=font, font_size=font_size, bold=True,
                 color=palette["dark_text"])

    stats_top = Inches(0.3) + Inches(box_height) + Inches(0.6)

    col_w = Inches(4)
    positions = [Inches(0.5), Inches(4.7), Inches(8.9)]

    for i, stat in enumerate(stats[:3]):
        x = positions[i]
        add_text_box(slide, stat.get("value", ""),
                     x, stats_top, col_w, Inches(2.2),
                     font_name=font, font_size=64, bold=True,
                     color=palette["primary"], align=PP_ALIGN.CENTER)
        add_text_box(slide, stat.get("label", ""),
                     x, stats_top + Inches(2.3), col_w, Inches(0.7),
                     font_name=font, font_size=15,
                     color=palette["dark_text"], align=PP_ALIGN.CENTER)
        if stat.get("sublabel"):
            add_text_box(slide, stat.get("sublabel"),
                         x, stats_top + Inches(3.1), col_w, Inches(0.5),
                         font_name=font, font_size=13, italic=True,
                         color="888888", align=PP_ALIGN.CENTER)


def render_two_column_slide(prs, slide_data: dict, palette: dict, font: str, image_bytes: Optional[bytes] = None):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)

    headline    = slide_data.get("headline", "")
    left_points = slide_data.get("left", [])
    right_text  = slide_data.get("right", "")

    font_size, box_height, _ = headline_fit(headline, box_width_in=12.0, max_font=30)
    add_text_box(slide, headline,
                 Inches(0.6), Inches(0.3), Inches(12), Inches(box_height),
                 font_name=font, font_size=font_size, bold=True,
                 color=palette["dark_text"])

    y = Inches(0.3) + Inches(box_height) + Inches(0.3)
    for pt in left_points[:5]:
        add_text_box(slide, f"• {pt}",
                     Inches(0.6), y, Inches(5.8), Inches(0.7),
                     font_name=font, font_size=15,
                     color=palette["dark_text"])
        y += Inches(0.85)

    panel_left, panel_top = Inches(7.2), Inches(1.4)
    panel_w, panel_h       = Inches(5.7), Inches(5.5)

    if image_bytes:
        # Real photo replaces the flat color panel entirely.
        add_picture_fit(slide, image_bytes, panel_left, panel_top, panel_w, panel_h)
        if right_text:
            # Caption strip under the image
            add_text_box(slide, right_text,
                         panel_left, panel_top + panel_h + Inches(0.1),
                         panel_w, Inches(0.8),
                         font_name=font, font_size=12, italic=True,
                         color=palette["dark_text"])
    elif right_text:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, panel_left, panel_top, panel_w, panel_h)
        box.fill.solid()
        box.fill.fore_color.rgb = hex_to_rgb(palette["secondary"])
        box.line.fill.background()

        add_text_box(slide, right_text,
                     panel_left + Inches(0.3), panel_top + Inches(0.3),
                     panel_w - Inches(0.6), panel_h - Inches(0.6),
                     font_name=font, font_size=15,
                     color=palette["dark_text"])


def render_closing_slide(prs, slide_data: dict, palette: dict, font: str, image_bytes: Optional[bytes] = None):
    layout = prs.slide_layouts[6]
    slide  = prs.slides.add_slide(layout)
    set_slide_background(slide, palette["primary"])

    add_text_box(slide, slide_data.get("title", "Thank You"),
                 Inches(1), Inches(2.2), Inches(11), Inches(1.5),
                 font_name=font, font_size=44, bold=True,
                 color=palette["accent"], align=PP_ALIGN.CENTER)

    sources = slide_data.get("sources", [])
    if sources:
        src_text = "Sources: " + " · ".join(sources)
        add_text_box(slide, src_text,
                     Inches(1), Inches(5.8), Inches(11), Inches(0.6),
                     font_name=font, font_size=12, italic=True,
                     color=palette["secondary"], align=PP_ALIGN.CENTER)


SLIDE_RENDERERS = {
    "title":             render_title_slide,
    "executive_summary": render_executive_summary,
    "content":           render_content_slide,
    "bullets":           render_content_slide,
    "chart":             render_chart_slide,
    "big_stat":          render_big_stat_slide,
    "two_column":        render_two_column_slide,
    "closing":           render_closing_slide,
}

# Slide types that can meaningfully use a supporting photo
IMAGE_CAPABLE_TYPES = {"content", "bullets", "two_column"}


# ── Endpoint 1: Generate slide JSON ───────────────────────────────────────────
@router.post("/generate-presentation")
async def generate_presentation(request: GeneratePresentationRequest):
    """
    Takes a search query + retrieved document chunks from AI Search.
    Groq structures the content into a slide JSON array.
    Lovable renders each slide as an editable card.
    """
    try:
        context = "\n\n---\n\n".join(request.chunks[:12])

        slide_types_note = ""
        if request.enable_charts:
            slide_types_note += ' Include at least 1-2 "chart" slides where data exists.'
        image_instruction = ""
        if request.enable_images:
            image_instruction = (
                ' For "content" and "two_column" slides where a supporting photo would '
                'help (not for slides that are already data-heavy like charts/stats), add an '
                '"image_query" field: a short 2-4 word visual search phrase (e.g. "team meeting '
                'office", "warehouse logistics") — not a restatement of the headline text.'
            )

        system_prompt = f"""You are an expert presentation designer creating executive-level slides.
Tone: {request.tone}. Target audience: CEOs and senior managers.
Create {request.num_slides} slides maximum. Every slide must have a clear headline that tells the story.{slide_types_note}{image_instruction}

Rules:
- Vary slide types — do not repeat the same layout twice in a row
- Headline should be a complete insight, not just a topic. Keep headlines under 90 characters where possible.
- Use ONLY data found in the provided context — mark estimated data with "(est.)"
- Charts must have realistic data arrays matching the labels length
- Big stat slides work well for 2-4 key numbers

Respond ONLY with valid JSON, no markdown fences:
{{
  "title": "presentation title",
  "slides": [
    {{
      "index": 0,
      "type": "title",
      "title": "string",
      "subtitle": "string"
    }},
    {{
      "index": 1,
      "type": "executive_summary",
      "headline": "string — complete insight sentence",
      "bullets": ["string", "string", "string"]
    }},
    {{
      "index": 2,
      "type": "big_stat",
      "headline": "string",
      "stats": [
        {{"value": "23%", "label": "Revenue Growth", "sublabel": "vs last quarter"}},
        {{"value": "₹10Cr", "label": "Enterprise ARR", "sublabel": "(est.)"}},
        {{"value": "4.2%", "label": "Churn Rate", "sublabel": "down from 6.1%"}}
      ]
    }},
    {{
      "index": 3,
      "type": "chart",
      "headline": "string — insight about the chart",
      "chart": {{
        "type": "column",
        "labels": ["Q1", "Q2", "Q3", "Q4"],
        "datasets": [
          {{"label": "2024", "data": [82, 95, 103, 118]}},
          {{"label": "2025", "data": [101, 124, null, null]}}
        ]
      }},
      "insight": "one-sentence insight about this data"
    }},
    {{
      "index": 4,
      "type": "two_column",
      "headline": "string",
      "left": ["point 1", "point 2", "point 3"],
      "right": "supporting paragraph or key context",
      "image_query": "short visual search phrase (optional)"
    }},
    {{
      "index": 5,
      "type": "content",
      "headline": "string",
      "bullets": ["string", "string", "string", "string"],
      "image_query": "short visual search phrase (optional)"
    }},
    {{
      "index": 6,
      "type": "content",
      "headline": "Recommendations",
      "bullets": ["string", "string", "string"]
    }},
    {{
      "index": 7,
      "type": "closing",
      "title": "Thank You",
      "sources": ["document1.pdf", "document2.docx"]
    }}
  ]
}}"""

        user_prompt = f"""Create a {request.tone} presentation about: "{request.query}"

Source content from company documents:
{context}

Generate {request.num_slides} slides. Use only information from the source content above.
Respond with JSON only."""

        data = ai.chat_json(
            messages=[{"role": "user", "content": user_prompt}],
            system=system_prompt,
            max_tokens=3000,
            temperature=0.3,
        )

        return {
            "title":  data.get("title", request.query),
            "slides": data.get("slides", []),
            "query":  request.query,
        }

    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"AI returned invalid JSON: {str(e)}")
    except Exception as e:
        import traceback
        print(f"GENERATE-PRESENTATION ERROR: {str(e)}")
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Presentation generation failed: {str(e)}")


# ── Endpoint 2: Edit one slide ─────────────────────────────────────────────────
@router.post("/edit-slide")
async def edit_slide(request: EditSlideRequest):
    """
    User types an instruction for one slide.
    Groq returns an updated slide JSON with the same index and type.
    Lovable replaces just that card in the sequence.
    """
    try:
        system_prompt = """You are editing a single presentation slide based on a user instruction.
Return ONLY the updated slide JSON with the same structure. Keep the same index and type unless the user explicitly asks to change the type.
If the slide has an "image_query" field, keep it unless the instruction is about the image.
Respond with valid JSON only, no markdown."""

        user_prompt = f"""Current slide:
{json.dumps(request.slide, indent=2)}

User instruction: {request.instruction}

Context: {request.context}

Return the updated slide JSON only."""

        updated_slide = ai.chat_json(
            messages=[{"role": "user", "content": user_prompt}],
            system=system_prompt,
            max_tokens=1000,
            temperature=0.4,
        )
        return updated_slide

    except json.JSONDecodeError as e:
        raise HTTPException(status_code=500, detail=f"AI returned invalid JSON: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Slide edit failed: {str(e)}")


# ── Endpoint 3: Export PPTX ────────────────────────────────────────────────────
@router.post("/export-pptx")
async def export_pptx(request: ExportPptxRequest):
    """
    Takes the final slide JSON array (after all user edits).
    Renders each slide with python-pptx using the branding settings.
    Fetches Unsplash images (best-effort) for slides that requested one.
    Uploads PPTX to R2 and returns a download URL.
    """
    try:
        branding = request.branding or BrandingSettings()

        palette = PALETTES.get(branding.palette or "midnight_executive",
                               PALETTES["midnight_executive"]).copy()

        if branding.primary_color:
            palette["primary"] = branding.primary_color.lstrip("#")

        font = branding.font or "Calibri"

        sorted_slides = sorted(request.slides, key=lambda s: s.get("index", 0))

        # Pre-fetch all needed images up front (async), before sync rendering
        images_by_index = await prefetch_images(sorted_slides, request.enable_images)

        prs = Presentation()
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H

        for slide_data in sorted_slides:
            slide_type = slide_data.get("type", "content")
            renderer   = SLIDE_RENDERERS.get(slide_type, render_content_slide)
            idx        = slide_data.get("index")
            img        = images_by_index.get(idx) if slide_type in IMAGE_CAPABLE_TYPES else None
            renderer(prs, slide_data, palette, font, img)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        pptx_bytes = buf.getvalue()

        safe_title = request.title.replace(" ", "_")[:50]
        key        = f"presentations/{uuid.uuid4()}/{safe_title}.pptx"
        url        = upload_to_r2(pptx_bytes, key, "application/vnd.openxmlformats-officedocument.presentationml.presentation")

        print(f"[export-pptx] Uploaded: {key} ({len(pptx_bytes)/1024:.1f}KB, {len(images_by_index)} images embedded)")

        return {
            "success":     True,
            "pptx_url":    url,
            "slide_count": len(request.slides),
            "images_embedded": len(images_by_index),
            "size_kb":     round(len(pptx_bytes) / 1024, 1),
        }

    except Exception as e:
        import traceback
        print(f"EXPORT-PPTX ERROR: {str(e)}")
        print(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")
